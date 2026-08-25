from __future__ import annotations

import json
import re
import tempfile
from dataclasses import dataclass
from pathlib import Path

from lxml import etree
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt

from .domain.supply import supply_state
from .paths import AppPaths, ensure_workspace


_manifest_cache: dict | None = None
_manifest_cache_path: Path | None = None
_manifest_cache_mtime_ns: int | None = None


def _load_manifest(paths: AppPaths | None = None) -> dict:
    """Load the current asset manifest, reloading after an in-app save.

    The size-rules page saves its profile definitions while the desktop server
    is running.  Preview and PowerPoint rendering share this helper, so a
    process-lifetime cache made those edits look ineffective until restart.
    """
    global _manifest_cache, _manifest_cache_path, _manifest_cache_mtime_ns
    active_paths = paths or ensure_workspace()
    manifest_path = active_paths.workspace_config_dir / "asset_manifest.json"
    try:
        mtime_ns = manifest_path.stat().st_mtime_ns
    except OSError:
        mtime_ns = None
    if (
        _manifest_cache is None
        or _manifest_cache_path != manifest_path
        or _manifest_cache_mtime_ns != mtime_ns
    ):
        try:
            _manifest_cache = json.loads(manifest_path.read_text("utf-8"))
        except Exception:
            _manifest_cache = {}
        _manifest_cache_path = manifest_path
        _manifest_cache_mtime_ns = mtime_ns
    return _manifest_cache


# ── Brand colors ──────────────────────────────────────────────────────────────
DTM_NAVY      = RGBColor(0x1E, 0x27, 0x61)
DTM_GRAY      = RGBColor(0x55, 0x55, 0x55)
DTM_DARKTEXT  = RGBColor(0x1A, 0x1A, 0x1A)
DTM_ORANGE    = RGBColor(0xC8, 0x60, 0x00)
DTM_ORANGE_BG = RGBColor(0xFF, 0xF0, 0xD4)
DTM_ALT_BG    = RGBColor(0xF2, 0xF2, 0xF5)
DTM_RED       = RGBColor(0xB8, 0x3A, 0x3A)
# Color system: BLUE = NEW, ORANGE = REUSED
TAG_NEW       = RGBColor(0x1A, 0x6F, 0xC8)   # blue  — new / installed
TAG_REUSED    = RGBColor(0xC8, 0x60, 0x00)   # orange — reused / transferred
_WHITE        = RGBColor(0xFF, 0xFF, 0xFF)
_LIGHT_GRAY   = RGBColor(0xCC, 0xCC, 0xCC)
_PANEL_BG     = RGBColor(0xF5, 0xF6, 0xFA)
_REUSED_BG    = RGBColor(0xFF, 0xF3, 0xE8)   # warm tint for reused parts
_NEW_BG       = RGBColor(0xF0, 0xF6, 0xFF)   # cool tint for new parts

# Slide dimensions (from template: 13.333" × 7.5")
SLIDE_W_EMU = 12191695
SLIDE_H_EMU = 6858000
FOOTER_H    = Inches(0.42)   # sticky bottom bar height (consistent on all slides)

VIEWS = ["front", "side", "top", "rear"]

PALETTE_TOKENS: dict[str, list[str]] = {
    "single": ["red", "blue", "white", "amber", "green"],
    "duo": ["red-white", "blue-white", "red-amber", "blue-amber", "red-green",
            "amber-white", "green-white", "green-amber"],
    "trio": ["red-blue-white", "red-blue-amber", "red-amber-white",
             "red-green-white", "blue-amber-white", "blue-green-white"],
}

BAR_SIZES = {"front": (4.110, 0.130), "rear": (4.110, 0.130), "side": (0.708, 0.12)}
BAR_TOP_SIZES = {
    "bar_roof_top":          (0.64, 2.7600),
    "bar_interior-front":    (0.270, 2.53),
    "bar_interior-rear":     (0.250, 2.000),
}
EQUIP_SIZES = {
    "Push Bumper": {"front": (3.128, 2.870), "side": (0.373, 1.353), "top": (0.342, 1.434)},
}

# ── Manifest layout ───────────────────────────────────────────────────────────
MANIFEST_LIGHT_CATS = {"warning_light", "scene_light", "light_bar"}
MANIFEST_STRUCTURAL_PART_TYPES = {
    "arm_rest", "console", "docking_station", "equipment_tray", "gun_lock",
    "motion_attachment", "pedestal_mount", "pit_bar", "push_bumper", "rear_partition",
    "seat_cover", "special_face_plate", "wing_wraps", "wire_covers",
}

# These related systems should be read together rather than be scattered by
# alphabetical order across a customer-facing manifest page.  The picker and
# build preview keep their existing names/order; this is presentation-only.
MANIFEST_SYSTEM_SORT_ORDER = {
    "push_bumper": 0,
    "pit_bar": 1,
    "wing_wraps": 2,
    "wire_covers": 3,
    "siren_speaker": 0,
    "howler": 1,
}

MANIFEST_COL_HEADERS   = ["PART / SKU", "DETAILS / SALES DESCRIPTION", "QTY", "LOCATION", "SOURCE"]
MANIFEST_COL_WIDTHS_IN = [3.15, 4.10, 0.48, 3.30, 1.30]

MANIFEST_TABLE_LEFT  = Inches(0.5)
MANIFEST_TABLE_TOP   = Inches(1.07)
MANIFEST_TABLE_W     = sum(Inches(w) for w in MANIFEST_COL_WIDTHS_IN)
MANIFEST_HDR_ROW_H   = Inches(0.30)
MANIFEST_DATA_MIN_H  = Inches(0.31)

# Diagram pages retain a short in-context warning for a few missing items.  A
# dedicated exception page is clearer (and cannot clip) when a stress-test or
# an incomplete configuration has many legitimate render failures.
INLINE_RENDER_FAILURE_LIMIT = 4
RENDER_EXCEPTION_COL_HEADERS = ["PART", "LOCATION", "RENDERING ISSUE"]
RENDER_EXCEPTION_COL_WIDTHS_IN = [4.15, 2.35, 5.83]


@dataclass
class _ManifestEntry:
    """One rendered row in the customer-facing parts manifest."""
    raw: object
    location: str
    name: str
    manufacturer: str = ""
    part_number: str = ""
    quantity: object = ""
    description: str = ""
    detail: str = ""
    comment: str = ""
    is_sku: bool = False
    indent: int = 0


# ─────────────────────────────────────────────────────────────────────────────
# Small helpers
# ─────────────────────────────────────────────────────────────────────────────

def _is_reused(part) -> bool:
    """Compatibility name: true for any customer-supplied hardware."""
    return supply_state(part).is_customer_supplied


def _source_label(part) -> str:
    return supply_state(part).label


def _color_label(part) -> str:
    """Return color string only (no lens)."""
    return (getattr(part, "raw_color", "") or getattr(part, "color", "")).strip()


def _lens_label(part) -> str:
    """Return lens string only (no 'Lens:' prefix — the value already contains 'Lens')."""
    return getattr(part, "lens", "").strip()


def _quantity_label(part) -> str:
    """Return a clear quantity for a build-page component card.

    Older workbook rows use zero/blank as the implicit single-unit default.
    A card only exists for a rendered component, so make that default explicit
    rather than showing an ambiguous dash.
    """
    quantity = getattr(part, "quantity", "")
    try:
        numeric = float(quantity)
    except (TypeError, ValueError):
        return "QTY: 1"
    if numeric <= 0:
        return "QTY: 1"
    rendered = str(int(numeric)) if numeric.is_integer() else str(numeric)
    return f"QTY: {rendered}"


def _color_lens_label(part) -> str:
    """Legacy combined helper used by manifest table. Returns 'Color  /  Lens' string."""
    color = _color_label(part)
    lens  = _lens_label(part)
    parts = []
    if color:
        parts.append(color)
    if lens:
        parts.append(lens)
    return "  /  ".join(parts)


# ─────────────────────────────────────────────────────────────────────────────
# Logo
# ─────────────────────────────────────────────────────────────────────────────

_cropped_logo_cache: Path | None = None


def _find_logo(paths: AppPaths | None = None) -> Path | None:
    active = paths or ensure_workspace()
    candidates = [
        active.workspace_assets_dir / "dtm_logo.png",
        active.assets_dir            / "dtm_logo.png",
    ]
    for p in candidates:
        if p.exists():
            return p
    return None


def _get_cropped_logo(logo_path: Path) -> Path:
    global _cropped_logo_cache
    if _cropped_logo_cache and _cropped_logo_cache.exists():
        return _cropped_logo_cache
    try:
        from PIL import Image
        with Image.open(logo_path) as img:
            if img.mode == "RGBA":
                bbox = img.getbbox()
                if bbox and bbox != (0, 0, img.width, img.height):
                    cropped = img.crop(bbox)
                    tmp = tempfile.NamedTemporaryFile(suffix=".png", delete=False)
                    cropped.save(tmp.name, "PNG")
                    _cropped_logo_cache = Path(tmp.name)
                    return _cropped_logo_cache
    except Exception:
        pass
    return logo_path


def place_logo(slide, paths: AppPaths | None = None, cover: bool = False) -> None:
    """Place the DTM logo in the top-right header band.

    cover=True makes it slightly larger for the cover page.
    """
    logo_path = _find_logo(paths)
    if not logo_path:
        return
    use_path = _get_cropped_logo(logo_path)
    try:
        from PIL import Image
        with Image.open(use_path) as img:
            img_w, img_h = img.size
        ratio    = img_w / img_h
        logo_h   = Inches(0.86) if cover else Inches(0.72)
        logo_w   = int(logo_h * ratio)
        margin   = Inches(0.12)
        logo_top = Inches(0.05)
        logo_left = SLIDE_W_EMU - logo_w - margin
        slide.shapes.add_picture(str(use_path), logo_left, logo_top,
                                 width=logo_w, height=logo_h)
    except Exception:
        pass


def place_logo_bottom(slide, paths: AppPaths | None = None) -> None:
    """Place the DTM logo just above the footer bar (for side/top view slides)."""
    logo_path = _find_logo(paths)
    if not logo_path:
        return
    use_path = _get_cropped_logo(logo_path)
    try:
        from PIL import Image
        with Image.open(use_path) as img:
            img_w, img_h = img.size
        ratio     = img_w / img_h
        logo_h    = Inches(0.30)
        logo_w    = int(logo_h * ratio)
        margin_r  = Inches(0.18)
        margin_b  = Inches(0.06)
        logo_top  = SLIDE_H_EMU - FOOTER_H - logo_h - margin_b
        logo_left = SLIDE_W_EMU - logo_w - margin_r
        slide.shapes.add_picture(str(use_path), logo_left, logo_top,
                                 width=logo_w, height=logo_h)
    except Exception:
        pass


# ─────────────────────────────────────────────────────────────────────────────
# Header / footer update (for template slides — view + notes slides)
# ─────────────────────────────────────────────────────────────────────────────

def add_slide_footer_bar(slide, footer_text: str) -> None:
    """Remove the template's transparent TextBox 5/4 and replace with a visible navy bar.

    Called for every view slide and the notes slide. Manifest slides build their
    own footer bar directly inside _make_manifest_slide.
    """
    for name in ("TextBox 2", "TextBox 4", "TextBox 5"):
        shape = find_shape(slide, name)
        if shape:
            shape._element.getparent().remove(shape._element)

    ftr_top = SLIDE_H_EMU - FOOTER_H
    ftr = slide.shapes.add_textbox(0, ftr_top, SLIDE_W_EMU, FOOTER_H)
    ftr.fill.solid()
    ftr.fill.fore_color.rgb = DTM_NAVY
    tf = ftr.text_frame
    tf.margin_left = Inches(0.3)
    tf.margin_top  = Inches(0.10)
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text           = footer_text
    r.font.size      = Pt(10)
    r.font.bold      = True
    r.font.color.rgb = _WHITE


def update_slide_header_footer(slide, title: str, subtitle: str = "", footer: str = "") -> None:
    """Replace the template header shapes with a full-width navy header band.

    The template's TextBox 1 has no fill (transparent), so white text would be
    invisible against the white slide background.  We delete TextBox 1/2 and
    create a fresh navy header that matches the cover/manifest slide style.
    TextBox 4/5 (footer) are left for add_slide_footer_bar() to replace.
    """
    for name in ("TextBox 1", "TextBox 2"):
        shape = find_shape(slide, name)
        if shape:
            shape._element.getparent().remove(shape._element)

    HDR_H = Inches(0.62)
    hdr = slide.shapes.add_textbox(0, 0, SLIDE_W_EMU, HDR_H)
    hdr.name = "DTM_SLIDE_HEADER"  # avoid collision with add_slide_footer_bar deletions
    hdr.fill.solid()
    hdr.fill.fore_color.rgb = DTM_NAVY
    tf = hdr.text_frame
    tf.margin_left = Inches(0.3)
    tf.margin_top  = Inches(0.08)
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text           = title
    r.font.size      = Pt(11)
    r.font.bold      = True
    r.font.color.rgb = _WHITE
    if subtitle:
        p2 = tf.add_paragraph()
        r2 = p2.add_run()
        r2.text           = subtitle
        r2.font.size      = Pt(11)
        r2.font.color.rgb = RGBColor(0xCC, 0xCC, 0xFF)


# ─────────────────────────────────────────────────────────────────────────────
# Layout helpers
# ─────────────────────────────────────────────────────────────────────────────

def _textbox(slide, left, top, width, height, text, font_size=10,
             bold=False, color=DTM_DARKTEXT, italic=False,
             align=PP_ALIGN.LEFT, wrap=True):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf  = box.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    r   = p.add_run()
    r.text           = text
    r.font.size      = Pt(font_size)
    r.font.bold      = bold
    r.font.italic    = italic
    r.font.color.rgb = color
    return box


def _kv_block(slide, items: list[tuple[str, str]], left, top, width,
              line_h: float = 0.26) -> int:
    y = top
    for key, value in items:
        if not value:
            continue
        box = slide.shapes.add_textbox(left, y, width, Inches(line_h + 0.06))
        tf  = box.text_frame
        tf.word_wrap = True
        p   = tf.paragraphs[0]
        kr  = p.add_run()
        kr.text           = f"{key}:  "
        kr.font.size      = Pt(10)
        kr.font.bold      = True
        kr.font.color.rgb = DTM_NAVY
        vr  = p.add_run()
        vr.text           = value
        vr.font.size      = Pt(12)
        vr.font.color.rgb = DTM_DARKTEXT
        y += Inches(line_h)
    return y


def _divider(slide, y_emu: int, width_in: float = 12.4, margin_left: float = 0.35) -> None:
    bar = slide.shapes.add_textbox(Inches(margin_left), y_emu, Inches(width_in), Inches(0.025))
    bar.fill.solid()
    bar.fill.fore_color.rgb = _LIGHT_GRAY


def _add_border(shape, color: RGBColor = _LIGHT_GRAY, width_pt: float = 0.75) -> None:
    sp_pr = shape._element.find(qn("p:spPr"))
    if sp_pr is None:
        return
    for old in sp_pr.findall(qn("a:ln")):
        sp_pr.remove(old)
    ln  = etree.SubElement(sp_pr, qn("a:ln"))
    ln.set("w", str(int(width_pt * 12700)))
    sf  = etree.SubElement(ln,  qn("a:solidFill"))
    clr = etree.SubElement(sf,  qn("a:srgbClr"))
    clr.set("val", str(color))


def _remove_border(shape) -> None:
    sp_pr = shape._element.find(qn("p:spPr"))
    if sp_pr is None:
        return
    for old in sp_pr.findall(qn("a:ln")):
        sp_pr.remove(old)
    ln = etree.SubElement(sp_pr, qn("a:ln"))
    etree.SubElement(ln, qn("a:noFill"))


def _stripe_box(slide, left, top, width, height, color: RGBColor) -> None:
    box = slide.shapes.add_textbox(left, top, width, height)
    box.fill.solid()
    box.fill.fore_color.rgb = color
    _remove_border(box)


def _card_bg(slide, left, top, width, height, color: RGBColor) -> None:
    box = slide.shapes.add_textbox(left, top, width, height)
    box.fill.solid()
    box.fill.fore_color.rgb = color
    _remove_border(box)


# ─────────────────────────────────────────────────────────────────────────────
# Cover slide
# ─────────────────────────────────────────────────────────────────────────────

def _find_part(parts, *names) -> object | None:
    """Return the first part whose name matches any of the given names (case-insensitive)."""
    targets = {n.lower() for n in names}
    for p in parts:
        if getattr(p, "name", "").lower() in targets:
            return p
    return None


_KNOWN_VEHICLE_MAKES = frozenset({
    "Ford", "Chevrolet", "Dodge", "Ram", "GMC", "Jeep", "Toyota",
    "Nissan", "Kia", "Tesla", "Honda", "Hyundai", "Volkswagen",
})


def _vehicle_fields(vehicle: dict, fallback_year: object = "", fallback_model: object = "") -> tuple[str, str, str, str]:
    """Normalize vehicle year/make/model when legacy drafts store all three in MODEL."""
    year = str(vehicle.get("YEAR", "") or fallback_year or "").strip()
    make = str(vehicle.get("MAKE", "") or "").strip()
    model = str(vehicle.get("MODEL", "") or fallback_model or "").strip()
    sub_model = str(vehicle.get("SUB MODEL", "") or "").strip()

    # The project build year is authoritative. Legacy drafts sometimes keep
    # the vehicle's original year in MODEL, so remove *any* leading model year
    # before splitting its make/model fields.
    if re.match(r"^\d{4}\s+", model):
        model = re.sub(r"^\d{4}\s+", "", model, count=1)
    if not make and model:
        candidate = model.split(None, 1)[0]
        if candidate.casefold() in {item.casefold() for item in _KNOWN_VEHICLE_MAKES}:
            make = candidate
    if make and model.casefold().startswith(make.casefold() + " "):
        model = model[len(make):].strip()
    return year, make, model, sub_model


def _project_vehicle_fields(info: dict) -> tuple[str, str, str, str]:
    new_v = info.get("NewVehicle", {}) or {}
    exist_v = info.get("ExistingVehicle", {}) or {}
    source = dict(new_v if any(new_v.values()) else exist_v)
    if info.get("BuildYear", ""):
        source["YEAR"] = info["BuildYear"]
    return _vehicle_fields(
        source,
        fallback_year=info.get("BuildYear", ""),
        fallback_model=info.get("VehicleType", ""),
    )


def _build_unit_label(build_type: object, unit_id: object) -> str:
    """Use the customer-facing ``Build Type #Number`` label wherever possible."""
    kind = str(build_type or "").strip()
    raw_unit = str(unit_id or "").strip()
    match = re.search(r"(\d+)\s*$", raw_unit)
    if kind and match:
        return f"{kind} #{match.group(1)}"
    if kind and raw_unit:
        return f"{kind} · {raw_unit}"
    if raw_unit and match:
        return f"Unit #{match.group(1)}"
    return raw_unit


def fill_overview(slide, project) -> None:
    info    = project.info
    new_v   = info.get("NewVehicle",      {})
    exist_v = info.get("ExistingVehicle", {})
    parts   = project.parts

    agency     = info.get("Agency",    "—")
    build_type = info.get("BuildType", "")
    year, make, model, sub_model = _project_vehicle_fields(info)
    unit_id   = (new_v.get("UNIT ID", new_v.get("UNIT", ""))
                 or exist_v.get("UNIT ID", exist_v.get("UNIT", "")))
    vin       = new_v.get("VIN",       "") or exist_v.get("VIN",       "")
    quote_num = info.get("QuoteNumber",    "")
    sales_rep = info.get("SalesRep",       "")

    # Remove slot shapes and template text boxes (prevents double footer / white-on-white)
    for name in ("PROJECT_INFO_BLOCK", "PARTS_TABLE_SLOT",
                 "TextBox 1", "TextBox 2", "TextBox 4", "TextBox 5"):
        shape = find_shape(slide, name)
        if shape:
            shape._element.getparent().remove(shape._element)

    # ── Navy header bar ───────────────────────────────────────────────────────
    hdr_title = "FLEET VEHICLE SPECIFICATION PACKAGE"
    if build_type:
        hdr_title += f"   •   {build_type.upper()}"
    hdr = slide.shapes.add_textbox(0, 0, SLIDE_W_EMU, Inches(0.95))
    hdr.fill.solid()
    hdr.fill.fore_color.rgb = DTM_NAVY
    tf = hdr.text_frame
    tf.margin_left = Inches(0.3)
    tf.margin_top  = Inches(0.21)
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text           = hdr_title
    r.font.size      = Pt(18)
    r.font.bold      = True
    r.font.color.rgb = _WHITE

    # ── Sticky footer bar — prominent, full vehicle identity ──────────────────
    vehicle_line = " ".join(filter(None, [year, make, model, sub_model]))
    unit_str     = _build_unit_label(build_type, unit_id)
    footer_parts = list(filter(None, [agency, vehicle_line, unit_str, "DTM Fleet Service"]))
    footer_text  = "   •   ".join(footer_parts)

    ftr_top = SLIDE_H_EMU - FOOTER_H
    ftr = slide.shapes.add_textbox(0, ftr_top, SLIDE_W_EMU, FOOTER_H)
    ftr.fill.solid()
    ftr.fill.fore_color.rgb = DTM_NAVY
    tf2 = ftr.text_frame
    tf2.margin_left = Inches(0.3)
    tf2.margin_top  = Inches(0.10)
    p2 = tf2.paragraphs[0]
    r2 = p2.add_run()
    r2.text           = footer_text
    footer_font_size = 10 if len(footer_text) <= 100 else 8.5 if len(footer_text) <= 135 else 7.5
    r2.font.size      = Pt(footer_font_size)
    r2.font.bold      = True
    r2.font.color.rgb = _WHITE

    # ── H1: Agency name — dominant ────────────────────────────────────────────
    L = Inches(0.45)
    y = Inches(1.05)

    agency_len = len(str(agency or "").strip())
    if agency_len <= 32:
        agency_font, agency_height = 48, 0.82
    elif agency_len <= 52:
        agency_font, agency_height = 39, 1.08
    elif agency_len <= 76:
        agency_font, agency_height = 32, 1.25
    else:
        agency_font, agency_height = 27, 1.42
    _textbox(slide, L, y, Inches(12.0), Inches(agency_height),
             agency, font_size=agency_font, bold=True, color=DTM_NAVY)
    y += Inches(agency_height)

    # ── H2: Vehicle identity ──────────────────────────────────────────────────
    veh_display = " ".join(filter(None, [year, make, model, sub_model]))
    if unit_str:
        veh_display += f"   |   {unit_str}"
    _textbox(slide, L, y, Inches(12.0), Inches(0.46),
             veh_display, font_size=24, bold=True,
             color=RGBColor(0x2A, 0x35, 0x80))
    y += Inches(0.46)

    _divider(slide, y, width_in=12.3, margin_left=0.45)
    y += Inches(0.18)

    # ── Two-column block ──────────────────────────────────────────────────────
    col_top  = y
    LEFT_W   = Inches(5.6)
    RIGHT_X  = Inches(6.5)
    RIGHT_W  = Inches(5.85)

    # Left: Quote + Sales info
    _textbox(slide, L, col_top, LEFT_W, Inches(0.25),
             "QUOTE & SALES INFORMATION", font_size=10, bold=True, color=DTM_NAVY)
    left_y = col_top + Inches(0.25)
    left_y = _kv_block(slide, [
        ("Sales Rep",    sales_rep),
        ("Quote #",      quote_num),
    ], L, left_y, LEFT_W)

    # Right: Vehicle specs — NEW primary card (always) + EXISTING secondary card (orange, if data)
    card_h = Inches(1.75)

    new_vehicle_for_display = dict(new_v)
    if info.get("BuildYear", ""):
        new_vehicle_for_display["YEAR"] = info["BuildYear"]
    new_year, new_make, new_model, new_sub_model = _vehicle_fields(
        new_vehicle_for_display, fallback_year=info.get("BuildYear", ""),
    )
    new_model_str = " ".join(filter(None, [new_model, new_sub_model]))
    new_unit      = new_v.get("UNIT ID", new_v.get("UNIT", ""))
    new_vin       = new_v.get("VIN",     "")

    ex_year, ex_make, ex_model, ex_sub_model = _vehicle_fields(exist_v)
    ex_model_str = " ".join(filter(None, [ex_model, ex_sub_model]))
    ex_unit      = exist_v.get("UNIT ID", exist_v.get("UNIT", ""))
    ex_vin       = exist_v.get("VIN",     "")
    exist_has_data = any([ex_year, ex_make, ex_model_str, ex_unit, ex_vin])

    if exist_has_data:
        CARD_W  = int((RIGHT_W - Inches(0.12)) / 2)
        EXIST_X = RIGHT_X + CARD_W + Inches(0.12)
    else:
        CARD_W  = RIGHT_W
        EXIST_X = None

    # Primary card: NEW VEHICLE (blue border)
    veh_bg = slide.shapes.add_textbox(RIGHT_X, col_top - Inches(0.04), CARD_W, card_h)
    veh_bg.fill.solid()
    veh_bg.fill.fore_color.rgb = _PANEL_BG
    _add_border(veh_bg, TAG_NEW, 1.0)
    _textbox(slide, RIGHT_X + Inches(0.10), col_top, CARD_W - Inches(0.2), Inches(0.25),
             "NEW VEHICLE", font_size=10, bold=True, color=TAG_NEW)
    _kv_block(slide, [
        ("Year",    new_year      or "—"),
        ("Make",    new_make      or "—"),
        ("Model",   new_model_str or "—"),
        ("Build",   unit_str      or new_unit or "—"),
        ("VIN",     new_vin       or "—"),
    ], RIGHT_X + Inches(0.10), col_top + Inches(0.25), CARD_W - Inches(0.2))

    # Secondary card: EXISTING VEHICLE (orange theme, only if data present)
    if exist_has_data:
        ex_bg = slide.shapes.add_textbox(EXIST_X, col_top - Inches(0.04), CARD_W, card_h)
        ex_bg.fill.solid()
        ex_bg.fill.fore_color.rgb = DTM_ORANGE_BG
        _add_border(ex_bg, DTM_ORANGE, 1.0)
        _textbox(slide, EXIST_X + Inches(0.10), col_top, CARD_W - Inches(0.2), Inches(0.25),
                 "EXISTING VEHICLE", font_size=10, bold=True, color=DTM_ORANGE)
        _kv_block(slide, [
            ("Year",    ex_year      or "—"),
            ("Make",    ex_make      or "—"),
            ("Model",   ex_model_str or "—"),
            ("Build",   unit_str     or ex_unit or "—"),
            ("VIN",     ex_vin       or "—"),
        ], EXIST_X + Inches(0.10), col_top + Inches(0.25), CARD_W - Inches(0.2))

    # ── Stats / tiles row ─────────────────────────────────────────────────────
    tiles_top = col_top + card_h + Inches(0.18)

    reused_count = sum(1 for p in parts if _is_reused(p))
    lights_count = sum(
        getattr(p, "quantity", 1) or 1
        for p in parts
        if getattr(p, "render_kind", "") == "light"
        and "tracer" not in getattr(p, "name", "").lower()
    )

    light_brands = sorted({
        _customer_manufacturer(getattr(p, "manufacturer", ""))
        for p in parts
        if getattr(p, "category", "") in MANIFEST_LIGHT_CATS
        and _customer_manufacturer(getattr(p, "manufacturer", ""))
    })
    brands_str = ", ".join(light_brands) if light_brands else "—"

    # --- Cage + Tray combined ---
    cage_part  = _find_part(parts, "front partition", "partition")
    cage_value = (getattr(cage_part, "part_number", "") or _customer_manufacturer(getattr(cage_part, "manufacturer", "")) or "Configured") if cage_part else "None"
    equip_part = _find_part(parts, "equipment tray")
    if equip_part:
        equip_value = (
            _customer_manufacturer(getattr(equip_part, "manufacturer", ""))
            or getattr(equip_part, "part_number", "")
            or getattr(equip_part, "name", "")
            or "Configured"
        )
    else:
        equip_value = "—"

    # --- Lighting System + Camera combined ---
    light_ctrl = _find_part(parts, "light controller", "lights controller")
    if light_ctrl:
        lc_lines = list(filter(None, [
            _customer_manufacturer(getattr(light_ctrl, "manufacturer", "")),
            getattr(light_ctrl, "part_number",  ""),
        ]))
        lighting_value = " · ".join(lc_lines) if lc_lines else "—"
    else:
        lighting_value = "—"
    camera_part = _find_part(parts, "camera dvr", "dvr", "camera system")
    if camera_part:
        cam_lines = list(filter(None, [
            _customer_manufacturer(getattr(camera_part, "manufacturer", "")),
            getattr(camera_part, "part_number",  ""),
        ]))
        camera_value = " · ".join(cam_lines) if cam_lines else "—"
    else:
        camera_value = ""

    # --- Bumper ---
    _BUMPER_KEYWORDS = ["push bumper", "pit bar", "wing wrap", "wire cover"]
    bumper_parts = [p for p in parts
                    if any(kw in getattr(p, "name", "").lower() for kw in _BUMPER_KEYWORDS)]
    bumper_mfg   = next((_customer_manufacturer(getattr(p, "manufacturer", "")) for p in bumper_parts
                         if _customer_manufacturer(getattr(p, "manufacturer", ""))), "")

    def _has(keyword: str) -> bool:
        return any(keyword.lower() in getattr(p, "name", "").lower() for p in bumper_parts)

    # ── 5-tile row: [Lights] [Reused] [Bumper] [Cage+Tray] [Lighting+Camera] ─
    SLIDE_USABLE_W = SLIDE_W_EMU - L - Inches(0.45)
    N_TILES  = 5
    TILE_GAP = Inches(0.10)
    TILE_W   = (SLIDE_USABLE_W - TILE_GAP * (N_TILES - 1)) / N_TILES

    def _tile_bg(tx, ty, th, bg_color, border_color):
        tb = slide.shapes.add_textbox(tx, ty, TILE_W, th)
        tb.fill.solid()
        tb.fill.fore_color.rgb = bg_color
        _add_border(tb, border_color, 0.5)
        return tb

    def _tile_label(tf, label, color=DTM_NAVY):
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text           = label
        r.font.size      = Pt(10)
        r.font.bold      = True
        r.font.color.rgb = color
        return tf

    def _tile_value(tf, value, font_size=12, bold=False, color=None, italic=False):
        p = tf.add_paragraph()
        r = p.add_run()
        r.text           = value
        r.font.size      = Pt(font_size)
        r.font.bold      = bold
        r.font.italic    = italic
        r.font.color.rgb = color or DTM_GRAY

    # Tile 0: Light Heads count
    tx = L
    BASE_TILE_H = Inches(1.10)
    tb = _tile_bg(tx, tiles_top, BASE_TILE_H, _PANEL_BG, _LIGHT_GRAY)
    tf = tb.text_frame
    tf.word_wrap   = True
    tf.margin_left = Inches(0.10)
    tf.margin_top  = Inches(0.06)
    p = tf.paragraphs[0]
    nr = p.add_run()
    nr.text           = str(lights_count)
    nr.font.size      = Pt(28)
    nr.font.bold      = True
    nr.font.color.rgb = DTM_NAVY
    _tile_value(tf, "Light Heads", font_size=12, color=DTM_NAVY)
    if brands_str and brands_str != "—":
        _tile_value(tf, brands_str, font_size=10, italic=True)

    # Tile 1: Reused / Transfer count
    tx = L + (TILE_W + TILE_GAP)
    tb = _tile_bg(tx, tiles_top, BASE_TILE_H, _REUSED_BG, _LIGHT_GRAY)
    tf = tb.text_frame
    tf.word_wrap   = True
    tf.margin_left = Inches(0.10)
    tf.margin_top  = Inches(0.06)
    p = tf.paragraphs[0]
    nr = p.add_run()
    nr.text           = str(reused_count)
    nr.font.size      = Pt(28)
    nr.font.bold      = True
    nr.font.color.rgb = TAG_REUSED
    _tile_value(tf, "Reused / Transfer", font_size=12, color=DTM_NAVY)

    # Tile 2: Bumper (inline checklist)
    tx = L + 2 * (TILE_W + TILE_GAP)
    tb = _tile_bg(tx, tiles_top, BASE_TILE_H, _PANEL_BG, _LIGHT_GRAY)
    tf = tb.text_frame
    tf.word_wrap   = True
    tf.margin_left = Inches(0.10)
    tf.margin_top  = Inches(0.06)
    _tile_label(tf, "Bumper")
    if not bumper_parts:
        _tile_value(tf, "✗  None", font_size=12, bold=True, color=DTM_RED)
    else:
        if bumper_mfg:
            _tile_value(tf, bumper_mfg, font_size=11)
        for label, keyword in [("Push bumper", "push bumper"), ("Pit bars", "pit bar"),
                               ("Wing wraps", "wing wrap"), ("Wire covers", "wire cover")]:
            if _has(keyword):
                _tile_value(tf, f"✓ {label}", font_size=8, color=DTM_GRAY)

    # Tile 3: Cage Type + Equipment Tray
    tx = L + 3 * (TILE_W + TILE_GAP)
    tb = _tile_bg(tx, tiles_top, BASE_TILE_H, _PANEL_BG, _LIGHT_GRAY)
    tf = tb.text_frame
    tf.word_wrap   = True
    tf.margin_left = Inches(0.10)
    tf.margin_top  = Inches(0.06)
    _tile_label(tf, "Cage / Tray")
    _tile_value(tf, f"Cage: {cage_value}", font_size=11)
    _tile_value(tf, f"Tray: {equip_value}", font_size=11)

    # Tile 4: Lighting System + Camera DVR
    tx = L + 4 * (TILE_W + TILE_GAP)
    tb = _tile_bg(tx, tiles_top, BASE_TILE_H, _PANEL_BG, _LIGHT_GRAY)
    tf = tb.text_frame
    tf.word_wrap   = True
    tf.margin_left = Inches(0.10)
    tf.margin_top  = Inches(0.06)
    _tile_label(tf, "Lighting / Camera")
    _tile_value(tf, f"Ctrl: {lighting_value}", font_size=11)
    if camera_value:
        _tile_value(tf, f"Cam: {camera_value}", font_size=11)


# ─────────────────────────────────────────────────────────────────────────────
# Parts manifest slides
# ─────────────────────────────────────────────────────────────────────────────

def _set_cell_bg(cell, color: RGBColor) -> None:
    tc   = cell._tc
    tcPr = tc.get_or_add_tcPr()
    for old in tcPr.findall(qn("a:solidFill")):
        tcPr.remove(old)
    sf  = etree.SubElement(tcPr, qn("a:solidFill"))
    clr = etree.SubElement(sf,   qn("a:srgbClr"))
    clr.set("val", str(color))


def _fmt_cell(cell, text: str, font_size: int = 8, bold: bool = False,
              color: RGBColor = DTM_DARKTEXT, bg: RGBColor | None = None,
              italic: bool = False, align=PP_ALIGN.LEFT) -> None:
    tf = cell.text_frame
    tf.word_wrap    = True
    tf.margin_left  = Inches(0.05)
    tf.margin_right = Inches(0.03)
    tf.margin_top   = Inches(0.025)
    tf.margin_bottom = Inches(0.025)
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.runs[0] if p.runs else p.add_run()
    r.text           = text
    r.font.size      = Pt(font_size)
    r.font.bold      = bold
    r.font.italic    = italic
    r.font.color.rgb = color
    if bg is not None:
        _set_cell_bg(cell, bg)


def _make_manifest_slide(
    prs,
    title: str,
    paths: AppPaths | None = None,
    footer_text: str = "",
    subtitle: str = "",
    title_font_size: int = 10,
):
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    for ph in list(slide.placeholders):
        ph._element.getparent().remove(ph._element)

    hdr = slide.shapes.add_textbox(0, 0, SLIDE_W_EMU, Inches(0.95))
    hdr.fill.solid()
    hdr.fill.fore_color.rgb = DTM_NAVY
    hdr.line.fill.background()

    # Keep category, context, and page number in independent text boxes.  A
    # single long header paragraph could wrap unpredictably on later pages in
    # PowerPoint-compatible renderers, leaving only the trailing page number.
    title_box = slide.shapes.add_textbox(Inches(0.3), Inches(0.09), Inches(10.1), Inches(0.34))
    tf = title_box.text_frame
    tf.clear()
    tf.word_wrap = False
    tf.margin_left = 0
    tf.margin_top  = 0
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text           = title
    r.font.size      = Pt(title_font_size)
    r.font.bold      = True
    r.font.color.rgb = _WHITE

    if subtitle:
        context, separator, page_number = subtitle.rpartition("   •   Page ")
        meta = slide.shapes.add_textbox(Inches(0.3), Inches(0.57), Inches(9.9), Inches(0.18))
        meta_tf = meta.text_frame
        meta_tf.clear()
        meta_tf.word_wrap = False
        meta_tf.margin_left = 0
        meta_tf.margin_top = 0
        meta_p = meta_tf.paragraphs[0]
        meta_r = meta_p.add_run()
        meta_r.text = context if separator else subtitle
        meta_r.font.size = Pt(9)
        meta_r.font.bold = True
        meta_r.font.color.rgb = RGBColor(0xD9, 0xDF, 0xF1)
        if separator:
            page = slide.shapes.add_textbox(Inches(10.35), Inches(0.57), Inches(2.55), Inches(0.18))
            page_tf = page.text_frame
            page_tf.clear()
            page_tf.word_wrap = False
            page_tf.margin_left = 0
            page_tf.margin_top = 0
            page_p = page_tf.paragraphs[0]
            page_p.alignment = PP_ALIGN.RIGHT
            page_r = page_p.add_run()
            page_r.text = f"Page {page_number}"
            page_r.font.size = Pt(9)
            page_r.font.bold = True
            page_r.font.color.rgb = RGBColor(0xD9, 0xDF, 0xF1)

    ftr_top = SLIDE_H_EMU - FOOTER_H
    ftr = slide.shapes.add_textbox(0, ftr_top, SLIDE_W_EMU, FOOTER_H)
    ftr.fill.solid()
    ftr.fill.fore_color.rgb = DTM_NAVY
    tf2 = ftr.text_frame
    tf2.margin_left = Inches(0.3)
    tf2.margin_top  = Inches(0.10)
    p2 = tf2.paragraphs[0]
    r2 = p2.add_run()
    r2.text           = footer_text or "DTM Fleet Service  •  Vehicle Build Sheet  •  Parts Manifest"
    r2.font.size      = Pt(10)
    r2.font.bold      = True
    r2.font.color.rgb = _WHITE

    place_logo(slide, paths)
    return slide


def _manifest_catalog(paths: AppPaths | None) -> dict[str, dict[str, dict[str, str]]]:
    """Index catalog descriptions without making an export depend on catalog I/O.

    A draft intentionally stores a durable SKU/model snapshot rather than an
    entire parts-db product.  The export can safely enrich that snapshot from
    the current local catalog when it is available, while old/offline builds
    remain exportable when it is not.
    """
    empty = {"sku": {}, "model": {}, "product": {}}
    active = paths or ensure_workspace()
    try:
        doc = json.loads(
            (active.workspace_config_dir / "parts_db.json").read_text("utf-8")
        )
    except Exception:
        return empty

    manufacturers = doc.get("manufacturers") or {}
    for product_id, product in (doc.get("products") or {}).items():
        if not isinstance(product, dict):
            continue
        manufacturer_id = str(product.get("manufacturer_id", "") or "")
        manufacturer = str(
            (manufacturers.get(manufacturer_id) or {}).get("label", manufacturer_id)
        ).strip()
        info = {
            "description": str(product.get("description", "") or "").strip(),
            "manufacturer": manufacturer,
            "model": str(product.get("model", "") or "").strip(),
            "color": "",
            "lens": "",
        }
        empty["product"][str(product_id)] = info
        if info["model"]:
            empty["model"][info["model"].casefold()] = info
        for sku in product.get("part_numbers") or []:
            if not isinstance(sku, dict):
                continue
            number = str(sku.get("part_number", "") or "").strip()
            if number:
                sku_info = dict(info)
                sku_info["description"] = str(
                    sku.get("qb_sales_description")
                    or sku.get("friendly_name")
                    or info["description"]
                    or ""
                ).strip()
                sku_info["color"] = " / ".join(
                    str(sku.get(field, "") or "").strip().title()
                    for field in ("color", "secondary_color", "tertiary_color")
                    if str(sku.get(field, "") or "").strip()
                )
                sku_info["lens"] = str(sku.get("lens_type", "") or "").strip()
                empty["sku"][number.casefold()] = sku_info
    return empty


def _picker_product_ids(value) -> list[str]:
    """Return catalog product ids retained anywhere in a picker snapshot."""
    found: list[str] = []
    if isinstance(value, dict):
        product_id = str(value.get("product_id", "") or "").strip()
        if product_id:
            found.append(product_id)
        for child in value.values():
            found.extend(_picker_product_ids(child))
    elif isinstance(value, list):
        for child in value:
            found.extend(_picker_product_ids(child))
    return found


def _catalog_detail(raw, part_number: str, catalog) -> tuple[str, str, bool, dict[str, str]]:
    """Return sales description, manufacturer fallback, SKU status, and SKU detail."""
    custom = (getattr(raw, "picker_config", {}) or {}).get("custom_part", {})
    if isinstance(custom, dict) and custom.get("description"):
        return str(custom["description"]).strip(), "", bool(part_number), {}

    normalized = str(part_number or "").strip().casefold()
    info = catalog["sku"].get(normalized)
    is_sku = info is not None
    if info is None:
        info = catalog["model"].get(normalized)
    if info is None:
        for product_id in _picker_product_ids(getattr(raw, "picker_config", {}) or {}):
            info = catalog["product"].get(product_id)
            if info is not None:
                break
    if not info:
        return "", "", False, {}
    return (
        info.get("description", ""),
        info.get("manufacturer", ""),
        is_sku,
        info,
    )


def _manifest_part_identity(entry: _ManifestEntry) -> tuple[str, str]:
    """Return a manifest row's clear product name and model/SKU support text."""
    prefix = "↳ " * entry.indent
    primary = f"{prefix}{_customer_visible_text(entry.name)}".strip()
    part_number = _customer_visible_text(entry.part_number)
    if part_number and entry.is_sku:
        part_number = f"{'SKUs' if ' / ' in part_number else 'SKU'}: {part_number}"
    secondary = "  ·  ".join(filter(None, [
        _customer_manufacturer(entry.manufacturer), part_number,
    ]))
    return primary, secondary


def _manifest_details(entry: _ManifestEntry) -> tuple[str, str, str]:
    """Return sales description, configured detail, and user comment."""
    return (
        _customer_visible_text(entry.description),
        _customer_visible_text(entry.detail),
        _customer_visible_text(entry.comment),
    )


def _manifest_line_count(text: str, width_inches: float) -> int:
    """Conservative line estimate used to keep table text inside its row."""
    if not text:
        return 1
    characters_per_line = max(9, int(width_inches * 11.5))
    total = 0
    for paragraph in str(text).splitlines() or [""]:
        words = paragraph.split()
        if not words:
            total += 1
            continue
        line_len = 0
        paragraph_lines = 1
        for word in words:
            word_len = len(word)
            if word_len > characters_per_line:
                if line_len:
                    paragraph_lines += 1
                    line_len = 0
                full, remainder = divmod(word_len, characters_per_line)
                paragraph_lines += max(0, full - (0 if remainder else 1))
                line_len = remainder
            elif not line_len:
                line_len = word_len
            elif line_len + 1 + word_len <= characters_per_line:
                line_len += 1 + word_len
            else:
                paragraph_lines += 1
                line_len = word_len
        total += paragraph_lines
    return total


def _manifest_row_height(entry: _ManifestEntry) -> int:
    """Size a compact row from its longest cell instead of letting text spill."""
    primary, secondary = _manifest_part_identity(entry)
    description, detail, comment = _manifest_details(entry)
    details = "\n".join(filter(None, [description, detail, f"Comment: {comment}" if comment else ""]))
    source = _source_label(entry.raw)
    line_count = max(
        _manifest_line_count("\n".join(filter(None, [primary, secondary])), MANIFEST_COL_WIDTHS_IN[0]),
        _manifest_line_count(details, MANIFEST_COL_WIDTHS_IN[1]),
        _manifest_line_count(str(entry.quantity or ""), MANIFEST_COL_WIDTHS_IN[2]),
        _manifest_line_count(entry.location, MANIFEST_COL_WIDTHS_IN[3]),
        _manifest_line_count(source, MANIFEST_COL_WIDTHS_IN[4]),
    )
    return max(MANIFEST_DATA_MIN_H, Inches(0.08 + 0.15 * line_count))


def _fmt_manifest_item_cell(cell, entry: _ManifestEntry, bg: RGBColor | None) -> None:
    primary, secondary = _manifest_part_identity(entry)
    tf = cell.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.03)
    tf.margin_top = Inches(0.025)
    tf.margin_bottom = Inches(0.025)
    p = tf.paragraphs[0]
    p.space_after = Pt(0)
    r = p.add_run()
    r.text = primary
    r.font.size = Pt(8 if entry.indent else 9)
    r.font.bold = not entry.indent
    r.font.color.rgb = DTM_NAVY if not entry.indent else DTM_DARKTEXT
    if secondary:
        p2 = tf.add_paragraph()
        p2.space_after = Pt(0)
        r2 = p2.add_run()
        r2.text = secondary
        r2.font.size = Pt(8)
        r2.font.color.rgb = DTM_GRAY
    if bg is not None:
        _set_cell_bg(cell, bg)


def _fmt_manifest_details_cell(cell, entry: _ManifestEntry, bg: RGBColor | None) -> None:
    description, detail, comment = _manifest_details(entry)
    tf = cell.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.03)
    tf.margin_top = Inches(0.025)
    tf.margin_bottom = Inches(0.025)
    if description:
        p = tf.paragraphs[0]
        p.space_after = Pt(0)
        r = p.add_run()
        r.text = description
        r.font.size = Pt(8)
        r.font.color.rgb = DTM_DARKTEXT
    if detail:
        p2 = tf.add_paragraph() if description else tf.paragraphs[0]
        p2.space_after = Pt(0)
        r2 = p2.add_run()
        r2.text = detail
        r2.font.size = Pt(8)
        r2.font.color.rgb = DTM_GRAY
    if comment:
        p3 = tf.add_paragraph() if (description or detail) else tf.paragraphs[0]
        p3.space_after = Pt(0)
        r3 = p3.add_run()
        r3.text = f"Comment: {comment}"
        r3.font.size = Pt(9)
        r3.font.bold = True
        r3.font.color.rgb = DTM_RED
    if not description and not detail and not comment:
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = "—"
        r.font.size = Pt(8)
        r.font.color.rgb = DTM_GRAY
    if bg is not None:
        _set_cell_bg(cell, bg)


def _part_row(table, row_idx: int, entry: _ManifestEntry, alt_bg: bool, row_h: int) -> None:
    state = supply_state(entry.raw)
    reused = state.is_customer_supplied
    used = reused and state.customer_condition == "used"
    row_bg = RGBColor(0xFF, 0xF3, 0xE8) if reused else (DTM_ALT_BG if alt_bg else None)

    source_text = _source_label(entry.raw)
    source_text = f"↺  {source_text}" if reused else f"■  {source_text}"
    source_color = DTM_RED if used else (TAG_REUSED if reused else TAG_NEW)
    source_bg = RGBColor(0xFF, 0xE4, 0xE4) if used else (
        RGBColor(0xFF, 0xE8, 0xD0) if reused else RGBColor(0xD8, 0xEC, 0xFF)
    )

    _fmt_manifest_item_cell(table.cell(row_idx, 0), entry, row_bg)
    _fmt_manifest_details_cell(table.cell(row_idx, 1), entry, row_bg)
    _fmt_cell(table.cell(row_idx, 2), str(entry.quantity or "—"),
              font_size=9, bold=True, bg=row_bg, align=PP_ALIGN.CENTER)
    _fmt_cell(table.cell(row_idx, 3), entry.location, font_size=9, bg=row_bg)
    _fmt_cell(table.cell(row_idx, 4), source_text, font_size=9 if used else 8,
              bold=used, color=source_color, bg=source_bg)
    table.rows[row_idx].height = row_h


def _clean_manifest_location(location: str, part_name: str = "") -> str:
    """Return a shop-useful location, omitting picker-only fixture identifiers."""
    text = str(location or "").strip()
    if not text or text.upper().startswith("FIXTURE:"):
        return ""

    location_words = set(re.findall(r"[a-z0-9]+", text.casefold()))
    part_words = set(re.findall(r"[a-z0-9]+", str(part_name or "").casefold()))
    # A location such as "Rear Interior Light Bar" adds no installation
    # information when the row itself is already named Rear Interior Light Bar.
    if location_words and (location_words == part_words or location_words <= part_words):
        return ""
    return text


def _clean_location(pp) -> str:
    raw_loc = (pp.raw.location or "").strip()
    part_name = str(getattr(pp.raw, "name", "") or "")
    if raw_loc:
        return _clean_manifest_location(raw_loc, part_name)
    if pp.placements:
        key = pp.placements[0].location_key or ""
        return _clean_manifest_location(key, part_name)
    return ""


def _manifest_group_for_part(pp) -> str:
    """Assign a manifest item to one of the customer-facing build categories."""
    raw = pp.raw
    category = (pp.category or "").lower()
    render_kind = (pp.render_kind or "").lower()
    part_type = (getattr(raw, "part_type", "") or "").lower()
    name = (getattr(raw, "name", "") or "").lower()

    if render_kind in {"light", "bar"} or category in MANIFEST_LIGHT_CATS:
        return "Lighting"
    if (
        part_type in MANIFEST_STRUCTURAL_PART_TYPES
        or any(keyword in name for keyword in (
            "bumper", "barrier", "cage", "console", "floor mat", "gun lock",
            "gunlock", "partition", "seat cover", "window bar",
        ))
    ):
        return "Structural Equipment"
    if category in {"custom", "unknown"} or render_kind in {"", "none"}:
        return "Other / Custom"
    return "Equipment & Electronics"


def _lighting_manifest_sort_rank(part) -> int:
    """Keep the lighting manifest in shop-friendly functional order."""
    part_type = (getattr(part, "part_type", "") or "").lower()
    name = (getattr(part, "name", "") or "").lower()
    combined = f"{part_type} {name}"
    # Older picker rows can retain ``warning_light`` as their generic type,
    # so identify tracers from the product name before the broader warning
    # classification below.
    if "tracer" in combined:
        return 4
    if "warning" in combined or part_type == "lighthead":
        return 0
    if "scene" in combined or "spotlight" in combined or "take-down" in combined:
        return 1
    if any(token in combined for token in ("interior", "dome", "rear_seat", "cargo")):
        return 2
    if "light_bar" in combined or part_type == "roof_light_bar":
        return 3
    return 5


def _manifest_entry_for_part(pp, catalog, *, indent: int = 0,
                             display_name: str | None = None) -> _ManifestEntry:
    raw = pp.raw
    part_number = str(getattr(raw, "part_number", "") or "").strip()
    description, catalog_manufacturer, is_sku, catalog_info = _catalog_detail(
        raw, part_number, catalog
    )
    visual = _manifest_visual_detail(raw, catalog_info=catalog_info)
    notes = _customer_manifest_notes(getattr(raw, "notes", ""))
    return _ManifestEntry(
        raw=raw,
        location=_customer_visible_text(_clean_location(pp)),
        name=_customer_manifest_name(raw, display_name),
        manufacturer=_customer_manufacturer(
            str(getattr(raw, "manufacturer", "") or "").strip() or catalog_manufacturer
        ),
        part_number=_customer_visible_text(part_number),
        quantity=getattr(raw, "quantity", ""),
        description=_customer_sales_description(description, raw),
        detail="\n".join(filter(None, [visual, notes])),
        comment=_customer_visible_text(getattr(raw, "comment", "")),
        is_sku=is_sku,
        indent=indent,
    )


_QB_IMPORT_REFERENCE = re.compile(
    r"\s*\(\s*(?:QB|QBO|QuickBooks)\s+Import\s*\)|\b(?:QB|QBO|QuickBooks)\s+Import\b",
    re.IGNORECASE,
)


def _customer_visible_text(value: object) -> str:
    """Strip accounting-source labels that do not belong in a build sheet."""
    cleaned_lines = []
    for line in str(value or "").splitlines():
        clean = _QB_IMPORT_REFERENCE.sub("", line)
        clean = re.sub(r"[ \t]{2,}", " ", clean).strip(" \t·-")
        if clean:
            cleaned_lines.append(clean)
    return "\n".join(cleaned_lines).strip()


def _customer_manufacturer(value: object) -> str:
    """Return a manufacturer label without internal QB-import provenance."""
    return _customer_visible_text(value)


def _customer_sales_description(value: object, raw: object | None = None) -> str:
    """Remove stale allocation quantities from customer-facing descriptions.

    The live quantity and location columns own those values. Some historical
    QBO descriptions contain a fixed ``(x2 cargo, x3 cage)`` recipe that is no
    longer true once round lights are allocated or edited in Builder.
    """
    allocation = (getattr(raw, "picker_config", {}) or {}).get("location_allocation")
    part_number = str(getattr(raw, "part_number", "") or "").strip().upper()
    round_light_description = (
        part_number in {"3SBCCDCR", "3SC0CDCR", "3SRCCDCR"}
        or '3" ROUND' in str(value or "").upper()
    )
    lines = []
    for line in str(value or "").splitlines():
        if (allocation or round_light_description) and re.match(
            r"^\s*\(\s*x\d+\b", line, re.IGNORECASE,
        ):
            continue
        clean = _customer_visible_text(line)
        if clean:
            lines.append(clean)
    return "\n".join(lines).strip()


def _customer_manifest_name(raw, display_name: str | None = None) -> str:
    """Return the short, customer-facing manifest name for a part.

    Sequence numbers help the editor and build preview distinguish placements,
    but they add no value in an order manifest.  In particular, a customer
    should see a quantity of siren speakers—not an implementation-specific
    ``Siren Speaker 1`` label.
    """
    name = str(display_name or getattr(raw, "name", "") or "").strip()
    part_type = str(getattr(raw, "part_type", "") or "").strip().casefold()
    if part_type == "siren_speaker" or re.fullmatch(r"siren speaker\s+\d+", name, re.I):
        return "Siren Speaker(s)"
    return _customer_visible_text(name)


def _customer_manifest_notes(value: object) -> str:
    """Remove picker-process language while preserving meaningful shop notes."""
    lines = str(value or "").splitlines()
    return "\n".join(
        _customer_visible_text(line)
        for line in lines
        if "guided system" not in line.casefold()
        and _customer_visible_text(line)
    ).strip()


def _customer_component_detail(value: object) -> str:
    """Keep useful configuration facts, not picker or shop-process narration."""
    detail = _customer_manifest_notes(value)
    normalized = detail.casefold().strip()
    if normalized in {
        "shop mounting location",
        "uses the selected center-console mic clip",
    }:
        return ""
    # A chosen location is already in the Location column; the rest of this
    # sentence described how the guided flow recorded it, not the part itself.
    if " — console position is set with the center console" in detail:
        detail = detail.split(" — console position is set with the center console", 1)[0].strip()
    return detail


def _is_included_control_head_component(parent: _ManifestEntry, component: dict) -> bool:
    """A PA microphone is included with the selected control head, not a line item."""
    part_type = str(getattr(parent.raw, "part_type", "") or "").strip().casefold()
    label = str(component.get("label", "") or component.get("name", "") or "").strip()
    return part_type == "control_head" and label.casefold() in {"pa mic", "pa microphone"}


def _manifest_system_sort_rank(raw) -> int:
    """Return a presentation rank for systems whose related pieces stay together."""
    part_type = str(getattr(raw, "part_type", "") or "").strip().casefold()
    if part_type in MANIFEST_SYSTEM_SORT_ORDER:
        return MANIFEST_SYSTEM_SORT_ORDER[part_type]

    # Old name-based builds do not have a part_type. Preserve the same
    # customer-facing grouping without changing the legacy records themselves.
    name = re.sub(r"\s+\d+$", "", str(getattr(raw, "name", "") or "").strip()).casefold()
    for prefix, rank in (
        ("push bumper", 0),
        ("pit bar", 1),
        ("wing wrap", 2),
        ("wire cover", 3),
        ("siren speaker", 0),
        ("howler", 1),
    ):
        if name.startswith(prefix):
            return rank
    return 99


def _manifest_visual_detail(raw, *, component: dict | None = None,
                            catalog_info: dict[str, str] | None = None,
                            include_parent_visual: bool = True) -> str:
    """Present explicit colour and lens details for the exact configured item."""
    component = component or {}
    catalog_info = catalog_info or {}
    tint = (getattr(raw, "picker_config", {}) or {}).get("window_tint")
    if isinstance(tint, dict):
        try:
            percentage = int(tint.get("percentage"))
        except (TypeError, ValueError):
            percentage = 0
        if percentage:
            return f"{percentage}% tint"
    color = (
        str(component.get("color", "") or "").strip()
        or (_color_label(raw) if include_parent_visual else "")
        or str(catalog_info.get("color", "") or "").strip()
    )
    lens = (
        str(component.get("lens", "") or component.get("lens_type", "") or "").strip()
        or (_lens_label(raw) if include_parent_visual else "")
        or str(catalog_info.get("lens", "") or "").strip()
    )
    # The catalog stores named lens variants (for example Smoked) explicitly.
    # Its older clear-lens lighting records intentionally leave that field
    # empty, so a colored light without a named variant is a clear lens.
    if not lens and color:
        lens = "Clear"
    details = [color] if color else []
    if lens:
        details.append(f"Lens: {lens.title()}")
    return "  ·  ".join(details)


def _combined_duo_manifest_entry(parent: _ManifestEntry, catalog) -> _ManifestEntry | None:
    """Collapse a standard driver/passenger DUO into one shop overview row.

    This is presentation-only. The draft's concrete SKU components remain
    separate for QuickBooks resolution and estimate quantities.
    """
    raw = parent.raw
    config = getattr(raw, "picker_config", {}) or {}
    duo_split = (
        config.get("colorsPerHead") == "duo" and config.get("mode") == "split"
    ) or bool(getattr(raw, "driver_color", "") and getattr(raw, "passenger_color", ""))
    if not duo_split:
        return None

    components = [
        component for component in (getattr(raw, "components", []) or [])
        if isinstance(component, dict) and str(component.get("part_number", "") or "").strip()
    ]
    if len(components) < 2:
        return None
    colors = [str(component.get("color", "") or "").strip() for component in components]
    if not any("red" in color.casefold() for color in colors) \
            or not any("blue" in color.casefold() for color in colors):
        return None

    def side_label(color: str) -> str:
        normalized = color.casefold()
        if "blue" in normalized and "red" not in normalized:
            return "passenger"
        if "red" in normalized and "blue" not in normalized:
            return "driver"
        return "configured"

    ordered = sorted(
        components,
        key=lambda component: 0 if side_label(str(component.get("color", ""))) == "passenger" else 1,
    )
    detail_parts = [
        f"{str(component.get('color', '') or '').strip()} ({side_label(str(component.get('color', '') or ''))})"
        for component in ordered if str(component.get("color", "") or "").strip()
    ]
    lens = _lens_label(raw) or "Clear"
    detail_parts.append(f"Lens: {lens.title()}")
    part_numbers = " / ".join(
        str(component.get("part_number", "") or "").strip() for component in ordered
    )
    quantity = 0
    for component in components:
        try:
            quantity += max(0, int(float(component.get("quantity", 1) or 0)))
        except (TypeError, ValueError):
            quantity += 1
    return _ManifestEntry(
        raw=raw,
        location=parent.location,
        name=parent.name,
        manufacturer=parent.manufacturer,
        part_number=part_numbers,
        quantity=quantity or parent.quantity,
        description=parent.description,
        detail="  ·  ".join(detail_parts),
        comment=parent.comment,
        is_sku=True,
        indent=parent.indent,
    )


def _manifest_component_entries(parent: _ManifestEntry, catalog, *,
                                collapse_skus: bool = False) -> list[_ManifestEntry]:
    """Expand SKU and guided-system details, optionally promoting SKUs to rows."""
    entries: list[_ManifestEntry] = []
    combined_duo = _combined_duo_manifest_entry(parent, catalog) if collapse_skus else None
    if combined_duo is not None:
        entries.append(combined_duo)
    promoted_sku_index = 0
    for component in getattr(parent.raw, "components", []) or []:
        if not isinstance(component, dict):
            continue
        if _is_included_control_head_component(parent, component):
            continue
        part_number = str(component.get("part_number", "") or "").strip()
        if combined_duo is not None and part_number:
            continue
        label = str(component.get("label", "") or "").strip()
        if not label and not part_number:
            continue
        if part_number:
            description, catalog_manufacturer, is_sku, catalog_info = _catalog_detail(
                parent.raw, part_number, catalog
            )
        else:
            # Guided-system components have no billable SKU of their own. The
            # parent kit's sales description belongs once on the parent row;
            # children keep only the installation detail selected by the user.
            description, catalog_manufacturer, is_sku, catalog_info = "", "", False, {}
        promote_sku = bool(part_number) and collapse_skus
        if promote_sku:
            row_name = parent.name
            row_indent = parent.indent
        elif label:
            row_name = label
            row_indent = parent.indent + 1
        else:
            row_name = "Selected SKU" if part_number else "Configured component"
            row_indent = parent.indent + 1
        visual_detail = _manifest_visual_detail(
            parent.raw,
            component=component,
            catalog_info=catalog_info,
            include_parent_visual=bool(part_number),
        )
        component_detail = _customer_component_detail(component.get("detail", ""))
        notes = ""
        if promote_sku and promoted_sku_index == 0:
            notes = _customer_manifest_notes(getattr(parent.raw, "notes", ""))
        component_has_supply = any(key in component for key in (
            "supply_type", "customer_condition", "customer_source",
            "new_or_used", "source",
        ))
        entries.append(_ManifestEntry(
            # Guided components can carry an independent supply decision. Use
            # it for their manifest status; older neutral SKU components keep
            # inheriting the parent status.
            raw=component if component_has_supply else parent.raw,
            location=_clean_manifest_location(
                component.get("location", ""), parent.name
            ) or parent.location,
            name=row_name,
            manufacturer=_customer_manufacturer(catalog_manufacturer or parent.manufacturer),
            part_number=_customer_visible_text(part_number),
            quantity=component.get("quantity", 1),
            description=_customer_sales_description(description, parent.raw),
            detail="\n".join(filter(None, [visual_detail, component_detail, notes])),
            comment=parent.comment if promote_sku and promoted_sku_index == 0 else "",
            is_sku=is_sku or bool(part_number),
            indent=row_indent,
        ))
        if promote_sku:
            promoted_sku_index += 1
    return entries


def _has_manifest_sku_components(raw) -> bool:
    return any(
        isinstance(component, dict)
        and str(component.get("part_number", "") or "").strip()
        for component in (getattr(raw, "components", []) or [])
    )


def _is_derived_manifest_part(candidate, planned_parts) -> bool:
    """Hide a render-only projection when its source component is already listed."""
    raw = candidate.raw
    if getattr(raw, "part_type", "") != "radio_antenna_top":
        return False
    line_id = str(getattr(raw, "line_id", "") or "")
    if not line_id:
        return False
    for other in planned_parts:
        if other is candidate:
            continue
        other_raw = other.raw
        if str(getattr(other_raw, "line_id", "") or "") != line_id:
            continue
        if any(
            isinstance(component, dict)
            and component.get("part_type") == "radio_antenna_top"
            for component in (getattr(other_raw, "components", []) or [])
        ):
            return True
    return False


def _manifest_groups(planned_parts, paths: AppPaths | None = None) -> list[tuple[str, list[_ManifestEntry]]]:
    """Build category pages while keeping a product's complete tree together."""
    catalog = _manifest_catalog(paths)
    visible_parts = [
        pp for pp in planned_parts
        if getattr(pp.raw, "include", True)
        and not is_render_only_part(pp)
        and not _is_derived_manifest_part(pp, planned_parts)
    ]
    by_line_id = {
        str(getattr(pp.raw, "line_id", "") or ""): pp
        for pp in visible_parts
        if getattr(pp.raw, "line_id", "")
    }
    children: dict[str, list] = {}
    roots: list = []
    for pp in visible_parts:
        parent_id = str(getattr(pp.raw, "parent_line_id", "") or "")
        if parent_id and parent_id in by_line_id and by_line_id[parent_id] is not pp:
            children.setdefault(parent_id, []).append(pp)
        else:
            roots.append(pp)

    def natural(value: str) -> tuple:
        return tuple(int(chunk) if chunk.isdigit() else chunk.casefold()
                     for chunk in re.split(r"(\d+)", value or ""))

    def sort_key(pp, label: str) -> tuple:
        raw = pp.raw
        return (
            _lighting_manifest_sort_rank(raw) if label == "Lighting" else 0,
            _manifest_system_sort_rank(raw),
            natural(str(getattr(raw, "name", "") or "")),
            natural(str(getattr(raw, "part_number", "") or "")),
            natural(_clean_location(pp)),
        )

    groups = {
        "Lighting": [],
        "Structural Equipment": [],
        "Equipment & Electronics": [],
        "Other / Custom": [],
    }
    visited: set[int] = set()

    def append_tree(pp, entries: list[_ManifestEntry], *, indent: int = 0,
                    display_name: str | None = None) -> None:
        marker = id(pp)
        if marker in visited:
            return
        visited.add(marker)
        entry = _manifest_entry_for_part(pp, catalog, indent=indent, display_name=display_name)
        has_sku_components = _has_manifest_sku_components(pp.raw)
        # Picker parents often hold a product/model label while their exact,
        # orderable SKU lives in ``components``.  The customer-facing manifest
        # should show the part name and that SKU together—not a generic parent
        # followed by a "Selected SKU" child row.  Preserve a direct parent
        # SKU as well when the parent is itself an orderable hardware item.
        if has_sku_components and not entry.is_sku:
            entries.extend(_manifest_component_entries(entry, catalog, collapse_skus=True))
        else:
            entries.append(entry)
            entries.extend(_manifest_component_entries(entry, catalog))
        parent_name = str(getattr(pp.raw, "name", "") or "")
        for child in sorted(children.get(str(getattr(pp.raw, "line_id", "") or ""), []),
                            key=lambda item: sort_key(item, _manifest_group_for_part(pp))):
            child_name = str(getattr(child.raw, "name", "") or "")
            prefix = f"{parent_name} · "
            if parent_name and child_name.startswith(prefix):
                child_name = child_name[len(prefix):]
            append_tree(child, entries, indent=indent + 1, display_name=child_name)

    for root in roots:
        label = _manifest_group_for_part(root)
        groups[label].append(root)

    result: list[tuple[str, list[_ManifestEntry]]] = []
    for label, roots_in_group in groups.items():
        entries: list[_ManifestEntry] = []
        for root in sorted(roots_in_group, key=lambda item: sort_key(item, label)):
            append_tree(root, entries)
        if entries:
            result.append((label, entries))
    return result


def is_render_only_part(part_or_planned) -> bool:
    raw = getattr(part_or_planned, "raw", part_or_planned)
    notes = (getattr(raw, "notes", "") or "").strip()
    part_number = (getattr(raw, "part_number", "") or "").strip()
    return notes.startswith("Included with Setina PB450L") or ":included-" in part_number


def add_parts_manifest_slides(prs, plan, paths: AppPaths | None = None) -> int:
    groups = _manifest_groups(plan.planned_parts, paths)
    if not groups:
        return 0

    proj    = plan.project
    new_v   = proj.get("NewVehicle",      {})
    exist_v = proj.get("ExistingVehicle", {})
    agency  = proj.get("Agency", "")
    year, make, model, sub_model = _project_vehicle_fields(proj)
    veh_line = " ".join(filter(None, [year, make, model, sub_model]))
    unit_id  = (new_v.get("UNIT ID", new_v.get("UNIT",""))
                or exist_v.get("UNIT ID", exist_v.get("UNIT","")))
    unit_part = _build_unit_label(proj.get("BuildType", ""), unit_id)
    hdr_parts = list(filter(None, [agency, veh_line, unit_part]))
    hdr_base  = "   •   ".join(hdr_parts)
    footer_text = "   •   ".join(filter(None, [agency, veh_line, unit_part, "DTM Fleet Service"]))

    n_cols     = len(MANIFEST_COL_HEADERS)
    col_widths = [Inches(w) for w in MANIFEST_COL_WIDTHS_IN]

    # Keep the full table inside the header/footer frame.  Row heights are
    # calculated from their contents rather than forcing text to spill below.
    avail_h = SLIDE_H_EMU - MANIFEST_TABLE_TOP - FOOTER_H - Inches(0.26)

    slides_added = 0
    page_num = 0

    for group_label, entries in groups:
        entry_index = 0
        category_page = 0

        while entry_index < len(entries):
            category_page += 1
            page_num += 1
            slide_rows: list[tuple] = []
            used_h = MANIFEST_HDR_ROW_H

            while entry_index < len(entries):
                entry = entries[entry_index]
                row_h = _manifest_row_height(entry)
                if used_h + row_h > avail_h and slide_rows:
                    break
                slide_rows.append((entry, row_h))
                used_h += row_h
                entry_index += 1

            category_suffix = "" if category_page == 1 else " — Continued"
            manifest_heading = f"{group_label.upper()}{category_suffix}"
            manifest_context = f"PARTS MANIFEST   •   {hdr_base}   •   Page {page_num}"
            slide = _make_manifest_slide(
                prs,
                manifest_heading,
                paths,
                footer_text=footer_text,
                subtitle=manifest_context,
                title_font_size=20,
            )
            slides_added += 1

            total_rows = 1 + len(slide_rows)
            tbl_shape = slide.shapes.add_table(
                total_rows, n_cols,
                MANIFEST_TABLE_LEFT, MANIFEST_TABLE_TOP,
                MANIFEST_TABLE_W, used_h,
            )
            table = tbl_shape.table
            for c, w in enumerate(col_widths):
                table.columns[c].width = w

            table.rows[0].height = MANIFEST_HDR_ROW_H
            for c, hdr in enumerate(MANIFEST_COL_HEADERS):
                _fmt_cell(table.cell(0, c), hdr,
                          font_size=9, bold=True, color=_WHITE,
                          bg=DTM_NAVY, align=PP_ALIGN.CENTER)

            for r_idx, (entry, row_h) in enumerate(slide_rows, start=1):
                _part_row(table, r_idx, entry, alt_bg=(r_idx % 2 == 0), row_h=row_h)

    return slides_added


def _render_exception_reason(item, view: str) -> str:
    """Make a planner/rendering message readable in the customer-facing report."""
    reason = getattr(item, "notes", "") or "Rendering failed for this placement"
    return reason.replace(f"{view}:", "").strip()


def _render_exception_row_height(item, view: str) -> int:
    """Estimate a table row height conservatively so exception pages never clip."""
    reason = _render_exception_reason(item, view)
    line_count = max(
        _manifest_line_count(getattr(item, "name", "") or "", RENDER_EXCEPTION_COL_WIDTHS_IN[0]),
        _manifest_line_count(getattr(item, "location", "") or "?", RENDER_EXCEPTION_COL_WIDTHS_IN[1]),
        _manifest_line_count(reason, RENDER_EXCEPTION_COL_WIDTHS_IN[2]),
    )
    return max(Inches(0.36), Inches(0.08 + 0.14 * line_count))


def add_render_exception_slides(
    prs,
    failures_by_view: list[tuple[str, list]],
    paths: AppPaths | None = None,
    footer_text: str = "",
) -> int:
    """Append paginated detail pages for diagram components that failed to render.

    The diagram itself keeps a compact red count so the visual remains usable;
    this companion page preserves the full list and its diagnostic reason.
    """
    n_cols = len(RENDER_EXCEPTION_COL_HEADERS)
    col_widths = [Inches(w) for w in RENDER_EXCEPTION_COL_WIDTHS_IN]
    table_left = Inches(0.50)
    table_top = Inches(1.32)
    table_w = sum(col_widths)
    hdr_h = Inches(0.30)
    avail_h = SLIDE_H_EMU - table_top - FOOTER_H - Inches(0.16)
    slides_added = 0

    for view, failures in failures_by_view:
        item_index = 0
        page_num = 0
        while item_index < len(failures):
            page_num += 1
            page_rows: list[tuple] = []
            used_h = hdr_h

            while item_index < len(failures):
                item = failures[item_index]
                row_h = _render_exception_row_height(item, view)
                if used_h + row_h > avail_h and page_rows:
                    break
                page_rows.append((item, row_h))
                used_h += row_h
                item_index += 1

            continuation = "" if page_num == 1 else " — Continued"
            title = f"RENDERING EXCEPTIONS — {view.upper()}{continuation}"
            slide = _make_manifest_slide(
                prs,
                title,
                paths,
                footer_text=footer_text or "DTM Fleet Service  •  Vehicle Build Sheet  •  Rendering Exceptions",
            )
            slides_added += 1

            intro = slide.shapes.add_textbox(
                table_left, Inches(1.00), table_w, Inches(0.20)
            )
            intro_tf = intro.text_frame
            intro_tf.margin_left = 0
            intro_tf.margin_top = 0
            intro_p = intro_tf.paragraphs[0]
            intro_r = intro_p.add_run()
            intro_r.text = (
                "Configured for this view but not drawn because its asset, instances, or placement could not be resolved."
            )
            intro_r.font.size = Pt(8)
            intro_r.font.color.rgb = DTM_GRAY

            table_shape = slide.shapes.add_table(
                1 + len(page_rows), n_cols, table_left, table_top, table_w, used_h
            )
            table = table_shape.table
            for c, width in enumerate(col_widths):
                table.columns[c].width = width
            table.rows[0].height = hdr_h
            for c, header in enumerate(RENDER_EXCEPTION_COL_HEADERS):
                _fmt_cell(
                    table.cell(0, c), header, font_size=9, bold=True,
                    color=_WHITE, bg=DTM_NAVY, align=PP_ALIGN.CENTER,
                )

            for row_idx, (item, row_h) in enumerate(page_rows, start=1):
                row_bg = DTM_ALT_BG if row_idx % 2 == 0 else None
                _fmt_cell(
                    table.cell(row_idx, 0), getattr(item, "name", "") or "?",
                    font_size=9, bold=True, color=DTM_NAVY, bg=row_bg,
                )
                _fmt_cell(
                    table.cell(row_idx, 1), getattr(item, "location", "") or "?",
                    font_size=9, bg=row_bg,
                )
                _fmt_cell(
                    table.cell(row_idx, 2), _render_exception_reason(item, view),
                    font_size=8, color=DTM_RED, bg=row_bg,
                )
                table.rows[row_idx].height = row_h

    return slides_added


# ─────────────────────────────────────────────────────────────────────────────
# Icon sizing
# ─────────────────────────────────────────────────────────────────────────────

def get_icon_size(part, icon_type, orient, view, icon_path_str="", paths: AppPaths | None = None):
    if icon_type == "equipment":
        equip = EQUIP_SIZES.get(part.name, {})
        if view in equip:
            return equip[view]
        if not icon_path_str:
            return (1.0, 1.0)
        from PIL import Image as PILImage
        with PILImage.open(
            (paths or ensure_workspace()).workspace_assets_dir / icon_path_str
        ) as img:
            image_w, image_h = img.size
        ratio   = image_w / image_h
        max_dim = 1.0
        return (max_dim, max_dim / ratio) if ratio > 1 else (max_dim * ratio, max_dim)

    if icon_type == "bar":
        # The roof light bar is a generic graphic — keep the legacy fixed
        # envelopes so it renders consistently across all builds. Every other
        # bar (interior front/rear and anything added later) uses a real
        # product photo, so we preserve the image's uploaded aspect ratio.
        is_roof_bar = "roof" in icon_path_str.lower()
        if is_roof_bar:
            if view == "top":
                for key, size in BAR_TOP_SIZES.items():
                    if key.replace("_", "-") in icon_path_str.replace("_", "-"):
                        return size
                return (0.40, 2.20)
            if view in BAR_SIZES:
                return BAR_SIZES[view]
            return (2.80, 0.25)

        max_dim = 2.5 if view == "top" else 1.0
        if not icon_path_str:
            return (max_dim, max_dim)
        try:
            from PIL import Image as PILImage
            with PILImage.open(
                (paths or ensure_workspace()).workspace_assets_dir / icon_path_str
            ) as img:
                image_w, image_h = img.size
        except Exception:
            return (max_dim, max_dim)
        ratio = image_w / image_h
        return (max_dim, max_dim / ratio) if ratio > 1 else (max_dim * ratio, max_dim)

    size_class = getattr(part, "size_class", "sm")
    defs       = _load_manifest(paths).get("size_rule_definitions", {})
    class_def  = defs.get(size_class) or defs.get("sm", {})
    views_data = class_def.get("views", {})
    view_data  = views_data.get(view) or views_data.get("front")
    if view_data:
        w, h = float(view_data["w"]), float(view_data["h"])
        return (h, w) if orient == "v" else (w, h)
    return (0.244, 0.085) if orient == "h" else (0.085, 0.244)


def icon_size_in_inches(
    render_kind: str,
    part_name: str,
    size_class: str,
    orientation: str,
    asset_path: str,
    view: str,
    size_override: "dict | None" = None,
    paths: "AppPaths | None" = None,
) -> "tuple[float, float]":
    """Single source of truth for icon sizing — used by both PPTX renderer and preview service.

    Returns (width_inches, height_inches).  size_override, if present, takes
    precedence over all catalog/manifest lookup.
    """
    if size_override and "w" in size_override and "h" in size_override:
        w, h = float(size_override["w"]), float(size_override["h"])
        return (h, w) if orientation == "v" else (w, h)
    from types import SimpleNamespace
    part = SimpleNamespace(name=part_name, size_class=size_class)
    return get_icon_size(part, render_kind, orientation, view, asset_path, paths=paths)


# ─────────────────────────────────────────────────────────────────────────────
# Generic shape helpers
# ─────────────────────────────────────────────────────────────────────────────

def _lock_picture_position(picture) -> None:
    c_nv_pic_pr = picture._element.find(".//" + qn("p:cNvPicPr"))
    if c_nv_pic_pr is None:
        return
    locks = c_nv_pic_pr.find(qn("a:picLocks"))
    if locks is None:
        locks = etree.SubElement(c_nv_pic_pr, qn("a:picLocks"))
    locks.set("noMove",   "1")
    locks.set("noResize", "1")


def find_shape(slide, name):
    for shape in slide.shapes:
        if shape.name == name:
            return shape
    return None


def slot_geometry(shape):
    left, top, width, height = shape.left, shape.top, shape.width, shape.height
    shape._element.getparent().remove(shape._element)
    return left, top, width, height


# ─────────────────────────────────────────────────────────────────────────────
# Vehicle image
# ─────────────────────────────────────────────────────────────────────────────

def place_vehicle_image(slide, vehicle_type, view, max_right_emu=None):
    """Return (picture_shape | None, img_box | None, equip_scale).

    img_box is (L,T,W,H) in EMU.  equip_scale is (scale_w, scale_h) — the ratio
    of the constrained vehicle image size to the unconstrained size.  Equipment
    icons (Push Bumper, Wing Wraps, etc.) whose sizes were calibrated for the
    unconstrained image should multiply their dimensions by equip_scale so they
    remain proportionally correct when the vehicle image is made smaller to
    accommodate a legend panel.

    For the side view the vehicle is placed in the bottom ~3" of the slide so
    the legend grid can occupy the top portion.
    max_right_emu: if set, constrains the vehicle image's right edge.
    """
    slot = find_shape(slide, "VEHICLE_IMAGE_SLOT")
    if not slot:
        return None, None, (1.0, 1.0)
    slot_left, slot_top, slot_w, slot_h = slot_geometry(slot)

    # Side/top view: move vehicle to bottom so legend grid fits above.
    # Slot top matches GRID_BOTTOM in place_legend_grid (3.50").
    if view in ("side", "top"):
        slot_left = 0
        slot_top  = Inches(3.50)
        slot_w    = SLIDE_W_EMU
        slot_h    = SLIDE_H_EMU - Inches(3.50) - FOOTER_H

    # Remember unconstrained slot width for equipment scale factor
    unconstrained_slot_w = slot_w

    # Honour caller-supplied right edge limit (e.g. to leave room for a legend panel)
    if max_right_emu is not None:
        available_w = max_right_emu - slot_left
        if available_w > 0:
            slot_w = min(slot_w, available_w)

    png = ensure_workspace().workspace_assets_dir / "vehicles" / f"{vehicle_type}_{view}.png"
    if not png.exists():
        return None, (slot_left, slot_top, slot_w, slot_h), (1.0, 1.0)

    from PIL import Image as PILImage
    with PILImage.open(png) as img:
        img_w, img_h = img.size

    img_ratio = img_w / img_h

    # Compute unconstrained final size (reference that size_per_view was calibrated for)
    uncons_slot_ratio = unconstrained_slot_w / slot_h
    if img_ratio > uncons_slot_ratio:
        ref_final_w = float(unconstrained_slot_w)
        ref_final_h = unconstrained_slot_w / img_ratio
    else:
        ref_final_h = float(slot_h)
        ref_final_w = slot_h * img_ratio

    # Compute constrained final size
    slot_ratio = slot_w / slot_h
    if img_ratio > slot_ratio:
        final_w    = slot_w
        final_h    = int(slot_w / img_ratio)
        final_left = slot_left
        final_top  = slot_top + (slot_h - final_h) // 2
    else:
        final_h    = slot_h
        final_w    = int(slot_h * img_ratio)
        final_left = slot_left + (slot_w - final_w) // 2
        final_top  = slot_top

    # Scale factor so equipment icons stay proportional to the vehicle image
    equip_scale = (final_w / ref_final_w, final_h / ref_final_h)

    pic = slide.shapes.add_picture(str(png), final_left, final_top,
                                   width=final_w, height=final_h)
    _lock_picture_position(pic)
    return pic, (final_left, final_top, final_w, final_h), equip_scale


# ─────────────────────────────────────────────────────────────────────────────
# Legend shared helpers
# ─────────────────────────────────────────────────────────────────────────────

def _est_wrapped_lines(text: str, inner_w_in: float, pt_size: int = 10) -> int:
    """Conservative estimate of how many display lines `text` needs.

    Uses an empirical average character width for Calibri mixed-case text, then
    applies a 0.85 safety factor to account for wide characters and word-wrap
    boundaries that push partial words to the next line.
    """
    if not text:
        return 0
    avg_char_w = pt_size * 0.0060          # inches per character at this pt size
    chars_per_line = max(1, int(inner_w_in / avg_char_w * 0.85))
    return max(1, -(-len(text) // chars_per_line))  # ceiling division


def _badge_label(part) -> tuple[str, object]:
    """Return (display_text, color) for the status badge on a legend card."""
    state = supply_state(part)
    if state.is_customer_supplied:
        condition = (
            "NEW" if state.customer_condition == "new"
            else "USED" if state.customer_condition == "used"
            else "CONDITION NEEDED"
        )
        return f"↺ CUSTOMER SUPPLIED / {condition}", TAG_REUSED
    return "■ NEW", TAG_NEW


def _legend_callouts(part) -> list[str]:
    """Return human-authored/source warnings shown on visual part cards."""
    callouts: list[str] = []
    comment = str(getattr(part, "comment", "") or "").strip()
    if comment:
        callouts.append(f"NOTE: {comment}")
    state = supply_state(part)
    if state.is_customer_supplied and state.customer_condition == "used":
        callouts.append(f"USED SOURCE: {state.customer_source or 'SOURCE NEEDED'}")
    return callouts


def _legend_headline(part) -> str:
    """Return the bold headline text for a legend card.

    For most parts the headline is the location (what the user typed in the workbook).
    Exceptions that keep the part *name* as headline:
      - Fixture parts (location is empty or starts with "FIXTURE:")
      - Parts whose name contains "tracer" or "light bar"
    """
    name       = getattr(part, "name",     "") or ""
    name_lower = name.lower()
    loc        = getattr(part, "location", "") or ""
    if "tracer" in name_lower or "light bar" in name_lower:
        return name
    if not loc or loc.upper().startswith("FIXTURE:"):
        return name
    return loc


def _legend_accessories(part, accessory_map: dict) -> list:
    """Return accessories for a concrete card, with legacy name-key fallback."""
    line_id = str(getattr(part, "line_id", "") or "")
    if line_id and accessory_map.get(line_id):
        return accessory_map[line_id]
    return accessory_map.get(getattr(part, "name", ""), [])


def _card_content_height(part, inner_w_in: float,
                          acc: dict,
                          pad_top=None, pad_bot=None,
                          lh_name=None, lh_spec=None, lh_acc=None,
                          min_card=None) -> int:
    """Compute card EMU height for a single part, accounting for text wrapping.

    All Inches() arguments use module defaults when None.
    """
    pad_top  = pad_top  if pad_top  is not None else Inches(0.06)
    pad_bot  = pad_bot  if pad_bot  is not None else Inches(0.04)
    lh_name  = lh_name  if lh_name  is not None else Inches(0.175)
    lh_spec  = lh_spec  if lh_spec  is not None else Inches(0.175)
    lh_acc   = lh_acc   if lh_acc   is not None else Inches(0.165)
    min_card = min_card if min_card is not None else Inches(0.34)

    badge_text, _ = _badge_label(part)
    headline = _legend_headline(part) + f"  {badge_text}"
    h = pad_top + _est_wrapped_lines(headline, inner_w_in, 10) * lh_name

    mfg   = _customer_manufacturer(getattr(part, "manufacturer", ""))
    pnum  = _customer_visible_text(getattr(part, "part_number", ""))
    color = _color_label(part)
    specs = "  ·  ".join(filter(None, [mfg, pnum, color, _quantity_label(part)]))
    if specs:
        h += _est_wrapped_lines(specs, inner_w_in, 11) * lh_spec
    lens = _lens_label(part)
    if lens:
        h += _est_wrapped_lines(lens, inner_w_in, 11) * lh_spec

    for callout in _legend_callouts(part):
        h += _est_wrapped_lines(callout, inner_w_in, 10) * Inches(0.18)

    for acc_entry in _legend_accessories(part, acc):
        acc_name = acc_entry[0] if isinstance(acc_entry, tuple) else acc_entry
        acc_pnum = acc_entry[1] if isinstance(acc_entry, tuple) else ""
        acc_text = "+ " + acc_name + (f"  ·  {acc_pnum}" if acc_pnum else "")
        h += _est_wrapped_lines(acc_text, inner_w_in, 10) * lh_acc

    return max(min_card, h + pad_bot)


# ─────────────────────────────────────────────────────────────────────────────
# Legend — standard sidebar (front / top / rear views)
# ─────────────────────────────────────────────────────────────────────────────

def place_legend(slide, placed, unplaced, accessory_map: dict | None = None,
                 view: str = "", panel_left_emu=None,
                 detail_unplaced: bool = False) -> None:
    """Two-column legend panel for front/rear views.

    panel_left_emu: left edge of the legend area in EMU (caller sets this to match
    whatever right edge was passed to place_vehicle_image).  Defaults to 7.80".
    """
    slot = find_shape(slide, "LEGEND_SLOT")
    if slot:
        slot._element.getparent().remove(slot._element)

    PANEL_LEFT  = panel_left_emu if panel_left_emu is not None else Inches(7.80)
    PANEL_TOP   = Inches(0.66)   # just below the 0.62" header band
    PANEL_W     = SLIDE_W_EMU - PANEL_LEFT - Inches(0.10)
    PANEL_H     = SLIDE_H_EMU - PANEL_TOP - FOOTER_H - Inches(0.06)
    BOTTOM_EDGE = PANEL_TOP + PANEL_H

    N_COLS   = 2
    COL_GAP  = Inches(0.08)
    COL_W    = (PANEL_W - COL_GAP * (N_COLS - 1)) / N_COLS
    STRIPE_W = Inches(0.05)

    # ── Panel section header ──────────────────────────────────────────────────
    count     = len(placed)
    hdr_label = (f"  {count} INSTALLED COMPONENT{'S' if count != 1 else ''}"
                 if placed else "  COMPONENTS")
    hdr_h = Inches(0.28)
    hdr   = slide.shapes.add_textbox(PANEL_LEFT, PANEL_TOP, PANEL_W, hdr_h)
    hdr.fill.solid()
    hdr.fill.fore_color.rgb = DTM_NAVY
    tf = hdr.text_frame
    tf.margin_left = Inches(0.08)
    tf.margin_top  = Inches(0.05)
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text           = hdr_label
    r.font.size      = Pt(10)
    r.font.bold      = True
    r.font.color.rgb = _WHITE

    cards_top = PANEL_TOP + hdr_h + Inches(0.04)

    # ── Empty state ───────────────────────────────────────────────────────────
    if not placed and not unplaced:
        msg = "NO ROOF-MOUNTED COMPONENTS" if view.lower() == "top" else "NO COMPONENTS ON THIS VIEW"
        tb = slide.shapes.add_textbox(PANEL_LEFT + Inches(0.10), cards_top,
                                       PANEL_W - Inches(0.20), Inches(0.60))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = msg; r.font.size = Pt(12); r.font.italic = True
        r.font.color.rgb = DTM_GRAY
        return

    # ── Card geometry ─────────────────────────────────────────────────────────
    CARD_GAP = Inches(0.05)
    PAD_TOP  = Inches(0.06)
    MIN_CARD = Inches(0.34)

    acc = accessory_map or {}

    # Inner column width (in inches) needed by _card_content_height
    _inner_w_in = float(COL_W - Inches(0.05) - Inches(0.08)) / 914400  # EMU → inches

    card_heights = [
        _card_content_height(p, _inner_w_in, acc, min_card=MIN_CARD)
        for p in placed
    ]

    # Reserve exactly what the inline diagnostic needs.  Larger failure sets are
    # listed on a dedicated page rather than forcing text below the footer.
    failure_reserve = (
        Inches(0.54) if detail_unplaced
        else Inches(0.30) + len(unplaced) * Inches(0.20)
        if unplaced else Inches(0.06)
    )
    cards_bottom = BOTTOM_EDGE - failure_reserve
    avail_h = cards_bottom - cards_top

    # Distribute cards into 2 columns greedily
    col_fills  = [0, 0]
    col_assign = []
    for h in card_heights:
        # pick the column with less fill, prefer left on tie
        c = 0 if col_fills[0] <= col_fills[1] else 1
        col_assign.append(c)
        col_fills[c] += h + CARD_GAP

    # If tallest column overflows, scale all heights down uniformly
    max_fill = max(col_fills) if col_fills else 0
    if max_fill > avail_h and card_heights:
        scale        = avail_h / max_fill
        card_heights = [max(MIN_CARD, h * scale) for h in card_heights]
        # Recompute fills
        col_fills = [0, 0]
        for h, c in zip(card_heights, col_assign):
            col_fills[c] += h + CARD_GAP

    # Track current y per column
    col_y = [cards_top, cards_top]

    for part, card_h, col_i in zip(placed, card_heights, col_assign):
        y   = col_y[col_i]
        cx  = PANEL_LEFT + col_i * (COL_W + COL_GAP)
        if y + card_h > cards_bottom:
            break

        reused       = getattr(part, "is_reused", False)
        stripe_color = TAG_REUSED if reused else TAG_NEW
        bg_color     = _REUSED_BG if reused else _NEW_BG

        _stripe_box(slide, cx, y, STRIPE_W, card_h - Inches(0.02), stripe_color)
        _card_bg(slide, cx + STRIPE_W, y, COL_W - STRIPE_W, card_h - Inches(0.02), bg_color)

        inner_x = cx + STRIPE_W + Inches(0.05)
        inner_w = COL_W - STRIPE_W - Inches(0.08)
        tb = slide.shapes.add_textbox(inner_x, y + PAD_TOP,
                                       inner_w, card_h - PAD_TOP)
        tf = tb.text_frame
        tf.word_wrap   = True
        tf.margin_top  = 0
        tf.margin_left = 0

        # Line 1: Location (or part name for fixtures/bars/tracers)  ■ NEW / ↺ REUSED
        badge_text, badge_color = _badge_label(part)
        p  = tf.paragraphs[0]
        nr = p.add_run()
        nr.text           = _legend_headline(part)
        nr.font.size      = Pt(10)
        nr.font.bold      = True
        nr.font.color.rgb = DTM_NAVY
        badge = p.add_run()
        badge.text           = f"  {badge_text}"
        badge.font.size      = Pt(10)
        badge.font.bold      = True
        badge.font.color.rgb = badge_color

        # Line 2: mfg · part# · color · qty (no lens here)
        mfg   = _customer_manufacturer(getattr(part, "manufacturer", ""))
        pnum  = _customer_visible_text(getattr(part, "part_number", ""))
        color = _color_label(part)
        specs = "  ·  ".join(filter(None, [mfg, pnum, color, _quantity_label(part)]))
        if specs:
            p2 = tf.add_paragraph()
            r2 = p2.add_run()
            r2.text           = specs
            r2.font.size      = Pt(11)
            r2.font.color.rgb = DTM_GRAY

        # Line 3: lens on its own line (value already contains "Lens")
        lens = _lens_label(part)
        if lens:
            p3 = tf.add_paragraph()
            r3 = p3.add_run()
            r3.text           = lens
            r3.font.size      = Pt(11)
            r3.font.color.rgb = DTM_GRAY

        for callout in _legend_callouts(part):
            pc = tf.add_paragraph()
            rc = pc.add_run()
            rc.text = callout
            rc.font.size = Pt(10)
            rc.font.bold = True
            rc.font.color.rgb = DTM_RED

        # Line 4+: accessories with part numbers
        for acc_entry in _legend_accessories(part, acc):
            acc_name = acc_entry[0] if isinstance(acc_entry, tuple) else acc_entry
            acc_pnum = acc_entry[1] if isinstance(acc_entry, tuple) else ""
            pa = tf.add_paragraph()
            ra = pa.add_run()
            ra.text           = "+ " + acc_name + (f"  ·  {acc_pnum}" if acc_pnum else "")
            ra.font.size      = Pt(10)
            ra.font.italic    = True
            ra.font.color.rgb = DTM_GRAY

        col_y[col_i] += card_h + CARD_GAP

    # ── "Not shown" box spans full panel width ────────────────────────────────
    if unplaced:
        y_note  = max(col_y) + Inches(0.06)
        box_h   = Inches(0.48) if detail_unplaced else Inches(0.26) + len(unplaced) * Inches(0.20)
        box_h   = min(box_h, BOTTOM_EDGE - y_note - Inches(0.04))
        bordered = slide.shapes.add_textbox(PANEL_LEFT, y_note, PANEL_W, box_h)
        bordered.fill.solid()
        bordered.fill.fore_color.rgb = RGBColor(0xFF, 0xF3, 0xF0)
        _add_border(bordered, DTM_RED, 0.5)
        tf = bordered.text_frame
        tf.word_wrap   = True
        tf.margin_left = Inches(0.08)
        tf.margin_top  = Inches(0.04)
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text           = (
            f"{len(unplaced)} COMPONENTS NOT SHOWN — SEE RENDERING EXCEPTIONS — {view.upper()}"
            if detail_unplaced else "ADDITIONAL COMPONENTS NOT SHOWN"
        )
        r.font.size      = Pt(8 if detail_unplaced else 10)
        r.font.bold      = True
        r.font.color.rgb = DTM_RED
        if not detail_unplaced:
            for upart in unplaced:
                loc = getattr(upart, "location", "") or "?"
                p2  = tf.add_paragraph()
                r2  = p2.add_run()
                r2.text           = f"  {upart.name}  @  {loc}"
                r2.font.size      = Pt(11)
                r2.font.color.rgb = DTM_RED


# ─────────────────────────────────────────────────────────────────────────────
# Legend — grid layout for side view (vehicle at bottom, cards at top)
# ─────────────────────────────────────────────────────────────────────────────

def place_legend_grid(slide, placed, unplaced, accessory_map: dict | None = None,
                      view: str = "side", detail_unplaced: bool = False) -> None:
    """Render part cards in a 4-column grid occupying the top portion of the slide.

    Used for side and top view slides where the vehicle image is at the bottom.
    """
    slot = find_shape(slide, "LEGEND_SLOT")
    if slot:
        slot._element.getparent().remove(slot._element)

    GRID_LEFT    = Inches(0.3)
    GRID_TOP     = Inches(0.66)   # flush with 0.62" header band bottom
    GRID_BOTTOM  = Inches(3.50)   # vehicle image starts at 3.50"
    GRID_COLS    = 4
    COL_GAP      = Inches(0.10)
    COL_W        = (SLIDE_W_EMU - GRID_LEFT * 2 - COL_GAP * (GRID_COLS - 1)) / GRID_COLS
    STRIPE_W     = Inches(0.06)
    CARD_GAP_V   = Inches(0.07)

    view_label = view.upper() if view else "SIDE"
    # Section header spanning full width
    count     = len(placed)
    hdr_label = (f"  {count} INSTALLED COMPONENT{'S' if count != 1 else ''}"
                 if placed else f"  {view_label} VIEW COMPONENTS")
    hdr_h = Inches(0.34)
    hdr   = slide.shapes.add_textbox(GRID_LEFT, GRID_TOP,
                                      SLIDE_W_EMU - GRID_LEFT * 2, hdr_h)
    hdr.fill.solid()
    hdr.fill.fore_color.rgb = DTM_NAVY
    tf = hdr.text_frame
    tf.margin_left = Inches(0.10)
    tf.margin_top  = Inches(0.07)
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text           = hdr_label
    r.font.size      = Pt(10)
    r.font.bold      = True
    r.font.color.rgb = _WHITE

    card_area_top  = GRID_TOP + hdr_h + Inches(0.06)
    failure_reserve = (
        Inches(0.52) if detail_unplaced
        else Inches(0.30) + len(unplaced) * Inches(0.20)
        if unplaced else 0
    )
    cards_bottom = GRID_BOTTOM - failure_reserve
    avail_h        = cards_bottom - card_area_top
    acc = accessory_map or {}
    last_card_bottom = card_area_top

    if not placed and not unplaced:
        tb = slide.shapes.add_textbox(GRID_LEFT + Inches(0.12), card_area_top,
                                       Inches(8), Inches(0.70))
        tf = tb.text_frame
        p  = tf.paragraphs[0]
        r  = p.add_run()
        r.text = f"NO {view_label}-VIEW COMPONENTS SPECIFIED"
        r.font.size = Pt(12); r.font.italic = True; r.font.color.rgb = DTM_GRAY
    else:
        # ── Dynamic per-row card heights ──────────────────────────────────────
        MIN_CARD  = Inches(0.32)
        inner_w_in = float(COL_W - STRIPE_W - Inches(0.08)) / 914400  # EMU → inches

        per_card_h = [
            _card_content_height(p, inner_w_in, acc, min_card=MIN_CARD)
            for p in placed
        ]

        # Group into rows; row height = max card height in that row
        n_rows    = max(1, -(-len(placed) // GRID_COLS))
        row_heights = []
        for r_i in range(n_rows):
            start = r_i * GRID_COLS
            end   = min(start + GRID_COLS, len(per_card_h))
            row_heights.append(max(per_card_h[start:end]) if per_card_h[start:end] else MIN_CARD)

        total_h = sum(row_heights) + CARD_GAP_V * (n_rows - 1)
        if total_h > avail_h and row_heights:
            scale       = avail_h / total_h
            row_heights = [max(MIN_CARD, h * scale) for h in row_heights]

        # Row y-offsets
        row_tops = [card_area_top]
        for rh in row_heights[:-1]:
            row_tops.append(row_tops[-1] + rh + CARD_GAP_V)

        for i, part in enumerate(placed):
            col_i  = i % GRID_COLS
            row_i  = i // GRID_COLS
            if row_i >= len(row_heights):
                break
            cx     = GRID_LEFT + col_i * (COL_W + COL_GAP)
            cy     = row_tops[row_i]
            card_h = row_heights[row_i]

            if cy + card_h > cards_bottom:
                break
            last_card_bottom = cy + card_h

            badge_text, badge_color = _badge_label(part)
            stripe_color = badge_color
            bg_color     = _REUSED_BG if stripe_color == TAG_REUSED else _NEW_BG

            _stripe_box(slide, cx, cy, STRIPE_W, card_h - Inches(0.02), stripe_color)
            _card_bg(slide, cx + STRIPE_W, cy,
                     COL_W - STRIPE_W, card_h - Inches(0.02), bg_color)

            inner_x = cx + STRIPE_W + Inches(0.05)
            inner_w = COL_W - STRIPE_W - Inches(0.08)
            tb = slide.shapes.add_textbox(inner_x, cy + Inches(0.04),
                                           inner_w, card_h - Inches(0.05))
            tf = tb.text_frame
            tf.word_wrap   = True
            tf.margin_top  = 0
            tf.margin_left = 0

            # Line 1: Location (or part name for fixtures/bars/tracers)  ■ NEW / ↺ REUSED
            p  = tf.paragraphs[0]
            nr = p.add_run()
            nr.text = _legend_headline(part); nr.font.size = Pt(10); nr.font.bold = True
            nr.font.color.rgb = DTM_NAVY
            bdg = p.add_run()
            bdg.text = f"  {badge_text}"; bdg.font.size = Pt(10); bdg.font.bold = True
            bdg.font.color.rgb = badge_color

            # Line 2: mfg · part# · color · qty (no lens)
            mfg   = _customer_manufacturer(getattr(part, "manufacturer", ""))
            pnum  = _customer_visible_text(getattr(part, "part_number", ""))
            color = _color_label(part)
            specs = "  ·  ".join(filter(None, [mfg, pnum, color, _quantity_label(part)]))
            if specs:
                p2 = tf.add_paragraph()
                r2 = p2.add_run()
                r2.text = specs; r2.font.size = Pt(11); r2.font.color.rgb = DTM_GRAY

            # Line 3: lens on its own line (value already contains "Lens")
            lens = _lens_label(part)
            if lens:
                p3 = tf.add_paragraph()
                r3 = p3.add_run()
                r3.text = lens; r3.font.size = Pt(11); r3.font.color.rgb = DTM_GRAY

            for callout in _legend_callouts(part):
                pc = tf.add_paragraph()
                rc = pc.add_run()
                rc.text = callout; rc.font.size = Pt(10); rc.font.bold = True
                rc.font.color.rgb = DTM_RED

            # Lines 4+: accessories
            for acc_entry in _legend_accessories(part, acc):
                acc_name = acc_entry[0] if isinstance(acc_entry, tuple) else acc_entry
                acc_pnum = acc_entry[1] if isinstance(acc_entry, tuple) else ""
                pa = tf.add_paragraph()
                ra = pa.add_run()
                ra.text = "+ " + acc_name + (f"  ·  {acc_pnum}" if acc_pnum else "")
                ra.font.size = Pt(10); ra.font.italic = True; ra.font.color.rgb = DTM_GRAY

    if unplaced:
        # Place "not shown" note just below the lowest rendered card row
        y_note = last_card_bottom + Inches(0.08)
        box_h  = Inches(0.46) if detail_unplaced else Inches(0.26) + len(unplaced) * Inches(0.20)
        box_h  = min(box_h, GRID_BOTTOM - y_note - Inches(0.04))
        bordered = slide.shapes.add_textbox(GRID_LEFT, y_note,
                                             SLIDE_W_EMU - GRID_LEFT * 2, box_h)
        bordered.fill.solid()
        bordered.fill.fore_color.rgb = RGBColor(0xFF, 0xF3, 0xF0)
        _add_border(bordered, DTM_RED, 0.5)
        tf = bordered.text_frame
        tf.word_wrap   = True
        tf.margin_left = Inches(0.10)
        tf.margin_top  = Inches(0.04)
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text           = (
            f"{len(unplaced)} COMPONENTS NOT SHOWN — SEE RENDERING EXCEPTIONS — {view.upper()}"
            if detail_unplaced else "ADDITIONAL COMPONENTS NOT SHOWN ON DIAGRAM"
        )
        r.font.size      = Pt(8 if detail_unplaced else 10)
        r.font.bold      = True
        r.font.color.rgb = DTM_RED
        if not detail_unplaced:
            for upart in unplaced:
                loc = getattr(upart, "location", "") or "?"
                p2  = tf.add_paragraph()
                r2  = p2.add_run()
                r2.text           = f"  {upart.name}  @  {loc}"
                r2.font.size      = Pt(11)
                r2.font.color.rgb = DTM_RED


# ─────────────────────────────────────────────────────────────────────────────
# Specify-palette swatches
# ─────────────────────────────────────────────────────────────────────────────

def place_specify_palette(slide, category: str, img_box, y_offset_emu: int = 0):
    tokens = PALETTE_TOKENS.get(category, [])
    if not tokens:
        return 0

    left_px, top_px, width_px, height_px = img_box
    icon_w   = Inches(0.45)
    icon_h   = Inches(0.114)
    label_h  = Inches(0.13)
    col_w    = Inches(0.54)
    max_cols = 10
    row_h    = icon_h + label_h + Inches(0.04)
    pal_top  = top_px + height_px + Inches(0.18) + y_offset_emu
    pal_left = left_px

    hdr = slide.shapes.add_textbox(pal_left, pal_top - Inches(0.17), Inches(4), Inches(0.15))
    tf  = hdr.text_frame
    r   = tf.paragraphs[0].add_run()
    r.text           = f"Specify {category} — keep one, delete the rest"
    r.font.size      = Pt(12)
    r.font.italic    = True
    r.font.color.rgb = DTM_RED

    rendered = 0
    for token in tokens:
        png_path = ensure_workspace().workspace_assets_dir / "lights" / f"sm_{token}_h.png"
        if not png_path.exists():
            continue
        col = rendered % max_cols
        row = rendered // max_cols
        x   = pal_left + col * col_w
        y   = pal_top  + row * row_h
        slide.shapes.add_picture(str(png_path), x, y, width=icon_w, height=icon_h)
        lbl = slide.shapes.add_textbox(x, y + icon_h + Inches(0.01), col_w, label_h)
        tf  = lbl.text_frame
        p   = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r2  = p.add_run()
        r2.text           = token
        r2.font.size      = Pt(6)
        r2.font.color.rgb = DTM_GRAY
        rendered += 1

    rows_used = (rendered + max_cols - 1) // max_cols if rendered else 0
    return rows_used * row_h + Inches(0.18)


# ─────────────────────────────────────────────────────────────────────────────
# Notes slide — structured sections
# ─────────────────────────────────────────────────────────────────────────────

NOTES_CATEGORIES = [
    "PROJECT-WIDE NOTES",
    "INSTALLATION NOTES",
    "DELIVERY REQUIREMENTS",
]


def fill_notes(slide, notes: dict[str, list[str]]) -> None:
    slot = find_shape(slide, "NOTES_SLOT")
    if not slot:
        return
    left, top, width, height = slot_geometry(slot)

    y       = top
    sec_gap = Inches(0.14)
    label_h = Inches(0.26)
    line_h  = Inches(0.22)
    placeholder_color = RGBColor(0xAA, 0xAA, 0xAA)

    _placeholders = {
        "PROJECT-WIDE NOTES":      "No project-wide notes specified.",
        "INSTALLATION NOTES":      "No installation notes specified.",
        "DELIVERY REQUIREMENTS":   "No delivery requirements specified.",
    }

    for section in NOTES_CATEGORIES:
        if y + label_h > top + height - Inches(0.1):
            break

        lbl = slide.shapes.add_textbox(left, y, width, label_h)
        lbl.fill.solid()
        lbl.fill.fore_color.rgb = RGBColor(0xEC, 0xEE, 0xF6)
        tf  = lbl.text_frame
        tf.margin_left = Inches(0.12)
        tf.margin_top  = Inches(0.05)
        p   = tf.paragraphs[0]
        r   = p.add_run()
        r.text           = section
        r.font.size      = Pt(10)
        r.font.bold      = True
        r.font.color.rgb = DTM_NAVY
        y += label_h + Inches(0.06)

        section_notes = notes.get(section, []) if isinstance(notes, dict) else []
        if section_notes:
            for idx, note in enumerate(section_notes):
                note_width = width - Inches(0.2)
                note_lines = _manifest_line_count(
                    f"{idx + 1}.  {note}", float(note_width) / 914400,
                )
                note_h = max(line_h, Inches(0.06 + 0.18 * note_lines))
                if y + note_h > top + height - Inches(0.1):
                    break
                tb  = slide.shapes.add_textbox(left + Inches(0.2), y,
                                                note_width, note_h)
                tf2 = tb.text_frame
                tf2.word_wrap = True
                tf2.margin_top = 0
                tf2.margin_bottom = 0
                p2  = tf2.paragraphs[0]
                nr  = p2.add_run()
                nr.text           = f"{idx + 1}.  "
                nr.font.size      = Pt(11)
                nr.font.bold      = True
                nr.font.color.rgb = DTM_NAVY
                br  = p2.add_run()
                br.text           = note
                br.font.size      = Pt(11)
                br.font.color.rgb = DTM_DARKTEXT
                y += note_h
        else:
            tb  = slide.shapes.add_textbox(left + Inches(0.2), y,
                                            width - Inches(0.2), line_h)
            tf2 = tb.text_frame
            p2  = tf2.paragraphs[0]
            r2  = p2.add_run()
            r2.text           = _placeholders.get(section, "No notes specified.")
            r2.font.size      = Pt(10)
            r2.font.italic    = True
            r2.font.color.rgb = placeholder_color
            y += line_h

        y += sec_gap
