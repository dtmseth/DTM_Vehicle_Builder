from __future__ import annotations

import logging
import re
import shutil
from datetime import datetime
from pathlib import Path
from types import SimpleNamespace

from lxml import etree
from pptx import Presentation
from pptx.oxml.ns import qn
from pptx.util import Inches

from .domain.geometry import slot_relative_positions
from .paths import AppPaths, ensure_workspace
from .ppt_helpers import (
    INLINE_RENDER_FAILURE_LIMIT,
    VIEWS as _LEGACY_VIEWS,
    _is_reused,
    _build_unit_label,
    _project_vehicle_fields,
    _source_label,
    add_parts_manifest_slides,
    add_render_exception_slides,
    add_slide_footer_bar,
    fill_notes,
    fill_overview,
    icon_size_in_inches,
    is_render_only_part,
    place_legend,
    place_legend_grid,
    place_logo,
    place_logo_bottom,
    place_specify_palette,
    place_vehicle_image,
    update_slide_header_footer,
)

_log = logging.getLogger(__name__)


def _project_shim(plan) -> SimpleNamespace:
    parts = []
    for pp in plan.planned_parts:
        if not pp.raw.include or is_render_only_part(pp):
            continue
        raw = pp.raw
        parts.append(SimpleNamespace(
            name         = pp.part_name,
            manufacturer = raw.manufacturer,
            part_number  = raw.part_number,
            color        = raw.raw_color,
            quantity     = raw.quantity,
            lens         = raw.lens,
            new_or_used  = raw.new_or_used,
            source       = raw.source,
            location     = raw.location,
            category     = pp.category or "",
            render_kind  = pp.render_kind or "",
        ))
    return SimpleNamespace(info=plan.project, parts=parts, notes=plan.notes)


def _legend_item(part_name: str, location: str = "", notes: str = "",
                 accessories: list[str] | None = None,
                 manufacturer: str = "", part_number: str = "",
                 color: str = "", lens: str = "",
                 new_or_used: str = "", source: str = "",
                 is_reused: bool = False, line_id: str = "",
                 quantity: object = "") -> SimpleNamespace:
    return SimpleNamespace(
        name         = part_name,
        location     = location,
        notes        = notes,
        accessories  = accessories or [],
        manufacturer = manufacturer,
        part_number  = part_number,
        color        = color,
        lens         = lens,
        new_or_used  = new_or_used,
        source       = source,
        is_reused    = is_reused,
        line_id      = line_id,
        quantity     = quantity,
    )


def _placement_key(part, placement) -> tuple[str, str, str, str]:
    """Return a stable identity for one planned placement on one view."""
    line_id = getattr(placement, "line_id", "") or getattr(part.raw, "line_id", "")
    return (
        line_id or part.part_name,
        part.part_name,
        placement.view,
        placement.location_key,
    )


def _record_render_failure(
    failures: dict[tuple[str, str, str, str], list[str]],
    key: tuple[str, str, str, str],
    messages: list[str],
) -> None:
    """Record one placement failure, preserving distinct messages in order."""
    bucket = failures.setdefault(key, [])
    for message in messages:
        if message and message not in bucket:
            bucket.append(message)


def _render_failures_for_view(plan, view: str, render_failures: dict):
    """Return only components expected on *view* that did not render.

    Manifest-only lines and parts meant for other vehicle views are intentionally
    absent here.  The red callout is a rendering diagnostic, not a duplicate
    manifest: it is reserved for a missing view placement or a failed asset.
    """
    issue_map: dict[tuple[str, str, str], SimpleNamespace] = {}
    prefix = f"{view}:"

    def add_issue(part, location: str, messages: list[str]) -> None:
        raw = part.raw
        key = (getattr(raw, "line_id", "") or part.part_name, part.part_name, location)
        issue = issue_map.get(key)
        if issue is None:
            issue = _legend_item(part.part_name, location, "")
            issue_map[key] = issue
        existing = [text for text in issue.notes.split("; ") if text]
        for message in messages:
            if message and message not in existing:
                existing.append(message)
        issue.notes = "; ".join(existing)

    for part in plan.planned_parts:
        if is_render_only_part(part):
            continue

        # A planner warning names this exact view only when a part should have
        # received a placement there but could not (for example, missing
        # coordinates). Generic warnings intentionally stay out of diagram pages.
        view_warnings = [warning for warning in part.warnings if warning.startswith(prefix)]
        if view_warnings:
            add_issue(part, part.raw.location or "?", view_warnings)

        for placement in part.placements:
            if placement.view != view or placement.color_profile == "specify_palette":
                continue
            messages = render_failures.get(_placement_key(part, placement), [])
            if messages:
                add_issue(part, part.raw.location or placement.location_key or "?", messages)

    return list(issue_map.values())


def _position_from_anchor(anchor: dict, img_box):
    left, top, width, height = img_box
    units = anchor.get("units", "relative_image")
    if units == "relative_image":
        return left + int(anchor["x"] * width), top + int(anchor["y"] * height)
    if units == "image_inches":
        return left + Inches(anchor["x"]), top + Inches(anchor["y"])
    return left + int(anchor["x"] * width), top + int(anchor["y"] * height)


def _slot_positions(pattern: str, slot_count: int, base_cx, base_cy,
                    img_box, h_spacing, v_spacing=None, slot_roles=None):
    """EMU-space adapter over domain geometry slot_relative_positions."""
    left, top, width, height = img_box
    anchor_x = (base_cx - left) / width
    anchor_y = (base_cy - top) / height
    h_rel = h_spacing / width
    v_rel = (v_spacing / height) if v_spacing is not None else None
    norm = slot_relative_positions(
        pattern, slot_count, anchor_x, anchor_y, h_rel, v_rel, slot_roles
    )
    return [(left + int(nx * width), top + int(ny * height)) for nx, ny in norm]


def _apply_shape_transforms(shape, rotation: float, flip_h: bool, flip_v: bool) -> None:
    if not rotation and not flip_h and not flip_v:
        return
    sp_pr = shape._element.find(qn("p:spPr"))
    if sp_pr is None:
        return
    xfrm = sp_pr.find(qn("a:xfrm"))
    if xfrm is None:
        xfrm = etree.SubElement(sp_pr, qn("a:xfrm"))
    if rotation:
        xfrm.set("rot", str(int(rotation * 60000)))
    if flip_h:
        xfrm.set("flipH", "1")
    if flip_v:
        xfrm.set("flipV", "1")


def _vertical_mirror_slot_is_reflected(placement, instance_index: int) -> bool:
    """Whether a vertical-mirror slot is the reflected (non-anchor) counterpart.

    Geometry emits vertical pairs top then bottom. The physical anchor side keeps
    the configured rotation, while the paired counterpart receives its reflected
    rotation. A centered pair is deterministic: top is the base and bottom is
    reflected.
    """
    if placement.pattern != "vertical_mirror" or len(placement.instances) <= 1:
        return False
    anchor_y = float((placement.anchor or {}).get("y", 0.5))
    is_top_slot = instance_index % 2 == 0
    if abs(anchor_y - 0.5) < 0.001:
        return not is_top_slot
    return is_top_slot != (anchor_y < 0.5)


def _stack_parts_above_vehicle(slide, vehicle_pic_element, part_elements: list[tuple]) -> None:
    """Keep the vehicle artwork below every rendered part in layer order.

    The vehicle image is the fixed base of every diagram.  Negative placement
    layers (for example, bumper hardware) mean "below other parts", never
    "behind the vehicle".  Grouped placements are supplied as their group
    element so the whole physical assembly keeps its relative stacking.
    """
    if vehicle_pic_element is None or not part_elements:
        return
    sp_tree = slide.shapes._spTree
    if vehicle_pic_element.getparent() is not sp_tree:
        return

    present = [
        (element, layer, index)
        for index, (element, layer) in enumerate(part_elements)
        if element.getparent() is sp_tree
    ]
    if not present:
        return
    for element, _, _ in present:
        sp_tree.remove(element)

    try:
        insert_at = list(sp_tree).index(vehicle_pic_element) + 1
    except ValueError:
        return
    for element, _, _ in sorted(present, key=lambda item: (item[1], item[2])):
        sp_tree.insert(insert_at, element)
        insert_at += 1


def _group_shapes(slide, shapes: list):
    if len(shapes) < 2:
        return None
    sp_tree = slide.shapes._spTree

    lefts  = [s.left              for s in shapes]
    tops   = [s.top               for s in shapes]
    rights = [s.left + s.width    for s in shapes]
    bots   = [s.top  + s.height   for s in shapes]
    gx, gy = min(lefts), min(tops)
    gcx    = max(rights) - gx
    gcy    = max(bots)   - gy

    existing_ids = {int(el.get("id", 0)) for el in sp_tree.iter() if el.get("id") is not None}
    grp_id = max(existing_ids, default=0) + 1

    grp_sp = etree.SubElement(sp_tree, qn("p:grpSp"))

    nv = etree.SubElement(grp_sp, qn("p:nvGrpSpPr"))
    cp = etree.SubElement(nv, qn("p:cNvPr"))
    cp.set("id", str(grp_id)); cp.set("name", f"Group {grp_id}")
    etree.SubElement(nv, qn("p:cNvGrpSpPr"))
    etree.SubElement(nv, qn("p:nvPr"))

    gpr  = etree.SubElement(grp_sp, qn("p:grpSpPr"))
    xfrm = etree.SubElement(gpr,    qn("a:xfrm"))
    off  = etree.SubElement(xfrm,   qn("a:off"));  off.set("x",  str(gx)); off.set("y", str(gy))
    ext  = etree.SubElement(xfrm,   qn("a:ext")); ext.set("cx", str(gcx)); ext.set("cy", str(gcy))
    cho  = etree.SubElement(xfrm,   qn("a:chOff")); cho.set("x", str(gx)); cho.set("y", str(gy))
    che  = etree.SubElement(xfrm,   qn("a:chExt")); che.set("cx", str(gcx)); che.set("cy", str(gcy))

    for shape in shapes:
        elem = shape._element
        sp_tree.remove(elem)
        grp_sp.append(elem)
    return grp_sp


def _instance_icon_size(placement, part_size, instance, view: str,
                        paths: AppPaths,
                        equip_scale: tuple[float, float] = (1.0, 1.0)) -> tuple[float, float]:
    """Return (width_inches, height_inches) for one icon instance.

    equip_scale is applied to equipment render_kind icons so they remain
    proportional to the vehicle image when it has been constrained by a
    legend panel (the sizes in size_per_view were calibrated for the
    unconstrained vehicle image).
    """
    w, h = icon_size_in_inches(
        render_kind=placement.render_kind,
        part_name=part_size.name,
        size_class=part_size.size_class,
        orientation=instance.orientation,
        asset_path=instance.asset_path,
        view=view,
        size_override=placement.size_override,
        paths=paths,
    )
    scale = getattr(placement, "size_scale", 1.0)
    # All icons anchored to the vehicle image must scale with it.  equip_scale
    # is (1.0, 1.0) when the vehicle is unconstrained (side/top views), so this
    # is a no-op in those cases.  For front/rear views with a legend panel the
    # vehicle is made smaller; all size_per_view / BAR_SIZES values were
    # calibrated for the full vehicle image, so every render_kind needs the
    # same proportional scale-down.
    w *= equip_scale[0]
    h *= equip_scale[1]
    return w * scale, h * scale


def _build_accessory_map(plan) -> dict[str, list[tuple[str, str]]]:
    """Return child purchase lines grouped by their exact parent legend card.

    Picker-created accessories use ``parent_line_id``.  That is the concrete
    parent selected in a build, unlike ``accessory_of`` which is a part-type
    relationship and may apply to more than one product.  Prefer the concrete
    relationship so each build-page card carries only its own children.  Keep
    the older name-keyed mapping as a fallback for legacy workbook builds.
    """
    acc_map: dict[str, list[tuple[str, str]]] = {}
    visible_parts = [
        pp for pp in plan.planned_parts
        if not is_render_only_part(pp)
    ]
    by_line_id = {
        str(getattr(pp.raw, "line_id", "") or ""): pp
        for pp in visible_parts
        if getattr(pp.raw, "line_id", "")
    }

    for pp in visible_parts:
        if is_render_only_part(pp):
            continue

        raw = pp.raw
        parent_id = str(getattr(raw, "parent_line_id", "") or "")
        parent = by_line_id.get(parent_id)
        parent_key = ""
        child_name = pp.part_name

        if parent is not None and parent is not pp:
            parent_key = parent_id
            parent_name = parent.part_name
            prefix = f"{parent_name} · "
            if child_name.casefold().startswith(prefix.casefold()):
                child_name = child_name[len(prefix):].strip()
        elif pp.accessory_of:
            # Legacy plans used the parent display name in accessory_of.
            parent_key = pp.accessory_of

        if parent_key:
            pnum = (getattr(raw, "part_number", "") or "").strip()
            entry = (child_name, pnum)
            entries = acc_map.setdefault(parent_key, [])
            if entry not in entries:
                entries.append(entry)
    return acc_map


def _move_slides_to_position(prs, start_position: int, count: int) -> None:
    """Move the last `count` slides in prs to start at index `start_position`."""
    if count <= 0:
        return
    sldIdLst = prs.slides._sldIdLst
    total    = len(sldIdLst)
    to_move  = list(sldIdLst)[total - count:]
    for elem in to_move:
        sldIdLst.remove(elem)
    for idx, elem in enumerate(to_move):
        sldIdLst.insert(start_position + idx, elem)


def _load_vehicle_view_config(vehicle_type: str, paths) -> tuple[list[str], dict]:
    """Return (external_view_order, {view_name: view_config}) from vehicle_layouts config.

    Falls back to the legacy hardcoded order if the config is unavailable.
    """
    import json
    try:
        layouts_path = paths.workspace_config_dir / "vehicle_layouts.json"
        layouts = json.loads(layouts_path.read_text("utf-8"))
        vehicle = layouts.get("vehicles", {}).get(vehicle_type, {})
    except Exception:
        vehicle = {}

    view_map   = vehicle.get("views", {})
    view_order = vehicle.get("view_order", list(_LEGACY_VIEWS))
    external   = [v for v in view_order
                  if view_map.get(v, {}).get("category", "external") == "external"]
    return (external or list(_LEGACY_VIEWS)), view_map


# The PPTX template has exactly this many view slides (indices 1..N, then notes at N+1)
_TEMPLATE_VIEW_SLOTS = 4


def _safe_part(s: str) -> str:
    """Sanitize a project info value into a safe filename segment."""
    s = re.sub(r'[\s/\\]+', '_', s.strip())
    s = re.sub(r'[^\w\-]', '', s)
    s = re.sub(r'_+', '_', s)
    return s.strip('_') or 'Unknown'


def build_output_filename(project: dict) -> str:
    """Build a timestamped export filename from project info."""
    agency     = _safe_part(str(project.get("Agency", "") or "Agency"))
    build_type = _safe_part(str(project.get("BuildType", "") or ""))
    new_v      = project.get("NewVehicle") or {}
    old_v      = project.get("ExistingVehicle") or {}
    unit       = _safe_part(str(new_v.get("UNIT ID", "") or old_v.get("UNIT ID", "") or "Unit"))
    year       = _safe_part(str(project.get("BuildYear", "") or new_v.get("YEAR", "") or old_v.get("YEAR", "") or "Year"))
    now        = datetime.now()
    hour       = now.hour % 12 or 12
    ampm       = "AM" if now.hour < 12 else "PM"
    ts         = now.strftime(f"%b%d_%Y_{hour}-{now.strftime('%M-%S')}{ampm}")
    parts      = [p for p in [agency, build_type, unit, year] if p and p != "Unknown"]
    return "_".join(parts) + f"_Updated_{ts}.pptx"


def render_plan_to_ppt(plan, paths: AppPaths | None = None) -> Path:
    active_paths = paths or ensure_workspace()
    template     = active_paths.templates_dir / "build_sheet_template.pptx"
    out_path     = active_paths.workspace_output_dir / build_output_filename(plan.project)

    vehicle_type = plan.project.get("VehicleType", "PIU")
    external_views, view_map = _load_vehicle_view_config(vehicle_type, active_paths)

    shutil.copyfile(template, out_path)
    prs      = Presentation(out_path)
    overview = prs.slides[0]
    # Map the first N external views to the template's view slide slots
    view_slides = {
        v: prs.slides[i + 1]
        for i, v in enumerate(external_views[:_TEMPLATE_VIEW_SLOTS])
    }
    notes_slide = prs.slides[_TEMPLATE_VIEW_SLOTS + 1]

    # ── Shared project metadata ───────────────────────────────────────────────
    project_shim = _project_shim(plan)
    new_v      = plan.project.get("NewVehicle",      {})
    exist_v    = plan.project.get("ExistingVehicle", {})
    agency     = plan.project.get("Agency", "")
    build_type = plan.project.get("BuildType", "")
    year, make, model, sub_model = _project_vehicle_fields(plan.project)
    veh_line = " ".join(filter(None, [year, make, model, sub_model]))
    unit_id  = (new_v.get("UNIT ID", new_v.get("UNIT",""))
                or exist_v.get("UNIT ID", exist_v.get("UNIT","")))
    unit_str = _build_unit_label(build_type, unit_id)
    footer   = "   •   ".join(filter(None, [agency, veh_line, unit_str, "DTM Fleet Service"]))

    # ── Cover slide ───────────────────────────────────────────────────────────
    fill_overview(overview, project_shim)
    place_logo(overview, active_paths, cover=True)

    # ── View slides ───────────────────────────────────────────────────────────
    accessory_map = _build_accessory_map(plan)
    detailed_render_failures: list[tuple[str, list]] = []

    for view in external_views:
        slide = view_slides.get(view)
        if slide is None:
            continue  # more external views than template slots; skip extras

        view_cfg      = view_map.get(view, {})
        legend_layout = view_cfg.get("legend_layout", "standard")
        logo_position = view_cfg.get("logo_position", "top-right")
        view_label    = view_cfg.get("label", view.upper())

        update_slide_header_footer(
            slide,
            title    = f"{view_label.upper()} VIEW — {agency}",
            subtitle = "  |  ".join(filter(None, [build_type, veh_line, unit_str])),
        )
        add_slide_footer_bar(slide, footer)
        if logo_position == "bottom":
            place_logo_bottom(slide, active_paths)
        else:
            place_logo(slide, active_paths)

        # For standard-legend views (front/rear) the legend occupies the right
        # portion of the slide; constrain the vehicle image to leave that room.
        _LEGEND_LEFT = Inches(5.80)   # legend panel starts here for standard layout
        veh_max_right = _LEGEND_LEFT if legend_layout != "grid" else None
        vehicle_pic_shape, img_box, equip_scale = place_vehicle_image(
            slide, vehicle_type, view, max_right_emu=veh_max_right
        )
        vehicle_pic_element = vehicle_pic_shape._element if vehicle_pic_shape else None
        if img_box is None:
            place_legend(slide, [], [], accessory_map, panel_left_emu=_LEGEND_LEFT)
            continue

        used_centers: list   = []
        collision            = Inches(0.18)
        rendered_part_elements: list[tuple] = []  # (shape XML, layer)
        specify_palettes: list[str] = []
        rendered_placement_keys: set[tuple[str, str, str, str]] = set()
        render_failures: dict[tuple[str, str, str, str], list[str]] = {}

        for pp in plan.planned_parts:
            for placement in pp.placements:
                if placement.view != view or not placement.instances:
                    if placement.view == view and not placement.instances:
                        _record_render_failure(
                            render_failures,
                            _placement_key(pp, placement),
                            ["No render instances planned"],
                        )
                    continue

                if placement.color_profile == "specify_palette":
                    cat = placement.instances[0].color_token
                    if cat not in specify_palettes:
                        specify_palettes.append(cat)
                    continue

                base_cx, base_cy = _position_from_anchor(placement.anchor, img_box)
                part_size        = SimpleNamespace(name=placement.part_name,
                                                   size_class=placement.size_class)
                first_inst       = placement.instances[0]
                if placement.render_kind in ("equipment", "bar") and not first_inst.asset_path:
                    _record_render_failure(
                        render_failures,
                        _placement_key(pp, placement),
                        first_inst.warnings or [
                            f"No asset resolved for '{placement.part_name}' ({view})"
                        ],
                    )
                    continue

                size_w, _ = _instance_icon_size(placement, part_size, first_inst, view,
                                                active_paths, equip_scale)
                icon_w_emu = Inches(size_w)

                if placement.h_spacing is not None and placement.h_spacing > 0:
                    if placement.h_spacing_units == "icon_width":
                        h_spacing_emu = int(icon_w_emu * placement.h_spacing)
                    else:
                        h_spacing_emu = int(img_box[2] * placement.h_spacing)
                else:
                    h_spacing_emu = int(icon_w_emu + Inches(0.03))

                if placement.v_spacing is not None and placement.v_spacing > 0:
                    if placement.h_spacing_units == "icon_width":
                        _, size_h = _instance_icon_size(placement, part_size, first_inst,
                                                        view, active_paths, equip_scale)
                        v_spacing_emu = int(Inches(size_h) * placement.v_spacing)
                    else:
                        v_spacing_emu = int(img_box[3] * placement.v_spacing)
                else:
                    v_spacing_emu = h_spacing_emu

                if placement.slot_indices and placement.position_slot_count:
                    all_positions = _slot_positions(
                        placement.pattern, placement.position_slot_count,
                        base_cx, base_cy, img_box,
                        h_spacing_emu, v_spacing_emu,
                    )
                    positions = [all_positions[i] for i in placement.slot_indices if i < len(all_positions)]
                else:
                    positions = _slot_positions(
                        placement.pattern, len(placement.instances),
                        base_cx, base_cy, img_box,
                        h_spacing_emu, v_spacing_emu,
                        slot_roles=[inst.slot_role for inst in placement.instances],
                    )
                translate_dx = int(img_box[2] * (getattr(placement, "translate_dx", 0.0) or 0.0))
                translate_dy = int(img_box[3] * (getattr(placement, "translate_dy", 0.0) or 0.0))
                if translate_dx or translate_dy:
                    positions = [
                        (px + translate_dx, py + translate_dy)
                        for px, py in positions
                    ]

                placement_shapes: list = []

                for instance_index, (instance, (px, py)) in enumerate(
                    zip(placement.instances, positions)
                ):
                    if not instance.asset_path:
                        _record_render_failure(
                            render_failures,
                            _placement_key(pp, placement),
                            instance.warnings or [
                                f"No asset resolved for '{placement.part_name}' ({view})"
                            ],
                        )
                        continue
                    icon_path = active_paths.workspace_assets_dir / instance.asset_path
                    if not icon_path.exists():
                        _record_render_failure(
                            render_failures,
                            _placement_key(pp, placement),
                            instance.warnings or [
                                f"Missing asset ({view}): {instance.asset_path}"
                            ],
                        )
                        if "wing_wrap_elitexd" in instance.asset_path.lower():
                            _log.warning(
                                "Wing Wrap EliteXD render skipped; asset missing: %s",
                                icon_path,
                            )
                        continue

                    inst_w, inst_h = _instance_icon_size(placement, part_size, instance,
                                                         view, active_paths, equip_scale)
                    iw = Inches(inst_w)
                    ih = Inches(inst_h)

                    ax, ay = px, py
                    if (not placement.group_shapes
                            and placement.render_kind not in ("bar", "equipment")):
                        nudge = Inches(0.30)
                        for ox, oy in used_centers:
                            if abs(ax - ox) < collision and abs(ay - oy) < collision:
                                ax    += nudge
                                nudge += Inches(0.30)
                    used_centers.append((ax, ay))

                    pic = slide.shapes.add_picture(
                        str(icon_path), ax - iw // 2, ay - ih // 2,
                        width=iw, height=ih,
                    )
                    if "wing_wrap_elitexd" in instance.asset_path.lower():
                        _log.info(
                            "Wing Wrap EliteXD picture inserted: path=%s view=%s "
                            "left=%s top=%s width=%s height=%s",
                            icon_path,
                            view,
                            ax - iw // 2,
                            ay - ih // 2,
                            iw,
                            ih,
                        )

                    is_right_slot = instance.slot_role in ("driver", "positive_x")
                    is_outer_edge = placement.pattern == "outer_edge_pillars"
                    is_sym_pattern = placement.pattern in ("mirror", "horizontal") or is_outer_edge
                    # Rear-view side roles are semantic (driver/passenger),
                    # not screen direction. Outer Edge's right column is
                    # therefore derived from the fixed three-head geometry.
                    is_right = (
                        instance_index >= len(placement.instances) // 2
                        if is_outer_edge else is_sym_pattern and is_right_slot
                    )
                    is_mirrored = (
                        placement.flip_mirrored_h
                        and is_right
                    )
                    inst_flip_h = placement.flip_h ^ is_mirrored
                    rotation_mirrored = (
                        is_right
                        or _vertical_mirror_slot_is_reflected(placement, instance_index)
                    )
                    inst_rot = (
                        (360 - placement.rotation) % 360
                        if rotation_mirrored and placement.rotation != 0 else placement.rotation
                    )
                    _apply_shape_transforms(pic, inst_rot, inst_flip_h, placement.flip_v)

                    layer = getattr(placement, "layer", 0)
                    placement_shapes.append(pic)
                    rendered_placement_keys.add(_placement_key(pp, placement))

                if placement.group_shapes and len(placement_shapes) >= 2:
                    group_element = _group_shapes(slide, placement_shapes)
                    if group_element is not None:
                        rendered_part_elements.append((group_element, layer))
                else:
                    rendered_part_elements.extend(
                        (shape._element, layer) for shape in placement_shapes
                    )

        # ── Layer Z-sort: vehicle first, then all parts back-to-front ────────
        _stack_parts_above_vehicle(slide, vehicle_pic_element, rendered_part_elements)

        # ── Build legend items ────────────────────────────────────────────────
        placed_legend: list   = []
        seen_placed: set[tuple] = set()
        for pp in plan.planned_parts:
            if is_render_only_part(pp):
                continue
            for pl in pp.placements:
                if pl.view != view:
                    continue
                if pl.color_profile == "specify_palette":
                    continue
                key = (pp.part_name, pl.location_key)

                if _placement_key(pp, pl) in rendered_placement_keys:
                    if key in seen_placed:
                        continue
                    seen_placed.add(key)
                    raw         = pp.raw
                    reused      = _is_reused(raw)
                    line_id = str(getattr(raw, "line_id", "") or "")
                    accessories = accessory_map.get(line_id) or accessory_map.get(pp.part_name, [])
                    placed_legend.append(_legend_item(
                        pp.part_name, raw.location or pl.location_key, raw.notes, accessories,
                        manufacturer = raw.manufacturer,
                        part_number  = raw.part_number,
                        color        = raw.raw_color,
                        lens         = raw.lens,
                        new_or_used  = raw.new_or_used,
                        source       = raw.source or "",
                        is_reused    = reused,
                        line_id      = line_id,
                        quantity     = raw.quantity,
                    ))

        # Keep bumper accessories in the legend as individual cards.  They are
        # small in the rendered vehicle view and otherwise easy to miss, even
        # when they are also summarized under the push bumper.

        unplaced_legend = _render_failures_for_view(plan, view, render_failures)
        detail_unplaced = len(unplaced_legend) > INLINE_RENDER_FAILURE_LIMIT
        if detail_unplaced:
            detailed_render_failures.append((view, unplaced_legend))
        palette_offset  = 0
        for cat in specify_palettes:
            consumed = place_specify_palette(slide, cat, img_box,
                                             y_offset_emu=palette_offset)
            if consumed:
                palette_offset += consumed

        if legend_layout == "grid":
            place_legend_grid(
                slide, placed_legend, unplaced_legend, accessory_map,
                view=view, detail_unplaced=detail_unplaced,
            )
        else:
            place_legend(
                slide, placed_legend, unplaced_legend, accessory_map,
                view=view, panel_left_emu=_LEGEND_LEFT,
                detail_unplaced=detail_unplaced,
            )

    # Dense or incomplete builds can have too many genuine failures to fit in a
    # diagram sidebar. Keep the full, view-specific list in paginated report
    # pages instead of clipping it or mixing it with the main manifest.
    if detailed_render_failures:
        add_render_exception_slides(
            prs, detailed_render_failures, active_paths, footer_text=footer
        )

    # ── Notes slide ───────────────────────────────────────────────────────────
    update_slide_header_footer(
        notes_slide,
        title    = f"BUILD NOTES — {agency}",
        subtitle = "  |  ".join(filter(None, [build_type, veh_line, unit_str])),
    )
    add_slide_footer_bar(notes_slide, footer)
    place_logo(notes_slide, active_paths)
    fill_notes(notes_slide, plan.notes)

    # ── Parts manifest slides (appended, then moved to position 1) ────────────
    n_manifest = add_parts_manifest_slides(prs, plan, active_paths)
    if n_manifest:
        _move_slides_to_position(prs, start_position=1, count=n_manifest)

    prs.save(out_path)
    return out_path
