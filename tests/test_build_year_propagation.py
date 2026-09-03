"""Regression coverage for project build-year propagation.

The project build year is the canonical year for generated build sheets. An
individual vehicle may still have an older legacy ``year`` value, but that
value must not leak into a project-tagged output.
"""
from __future__ import annotations

from dataclasses import replace
from pathlib import Path
from types import SimpleNamespace

from dtm_buildsheet.app.services.draft_service import handle_generate_from_draft
from dtm_buildsheet.app.services.project_service import handle_create_individual_draft
from dtm_buildsheet.app.services.qb_sync_service import _job_display_name
from dtm_buildsheet.domain.plan_models import BuildPlan
from dtm_buildsheet.domain.project_models import BuildUnit, CustomerInfo, IndividualUnit
from dtm_buildsheet.inputs.project_drafts import DraftPart, load_draft, new_draft, save_draft
from dtm_buildsheet.inputs.project_entry import new_project, save_project
from dtm_buildsheet.paths import AppPaths
from dtm_buildsheet.render_ppt import build_output_filename, render_plan_to_ppt


def _paths(tmp_path: Path) -> AppPaths:
    return AppPaths(
        workspace_dir=tmp_path,
        workspace_drafts_dir=tmp_path / "drafts",
        workspace_projects_dir=tmp_path / "projects",
        workspace_output_dir=tmp_path / "output",
        workspace_config_dir=tmp_path / "config",
        workspace_presets_dir=tmp_path / "presets",
    )


def _granite_project(draft_id: str) -> object:
    individual = IndividualUnit(
        individual_id="ind-1",
        unit_number="02",
        year="2026",
        make="Ford",
        model="Police Interceptor Utility",
        draft_id=draft_id,
    )
    return new_project(
        project_id="granite-2027",
        customer=CustomerInfo(agency="Granite Falls Police Department", build_year="2027"),
        build_units=[
            BuildUnit(
                unit_id="unit-1",
                vehicle_model="PIU",
                build_type="Patrol",
                individuals=[individual],
            )
        ],
    )


def test_individual_draft_uses_project_build_year(tmp_path):
    paths = _paths(tmp_path)
    project = _granite_project("")
    save_project(project, paths)

    result = handle_create_individual_draft(
        project.project_id, "unit-1", "ind-1", paths
    )

    assert result["ok"] is True
    draft = load_draft(result["draft_id"], paths.workspace_drafts_dir)
    assert draft.vehicle_info["BuildYear"] == "2027"
    assert draft.vehicle_info["NewVehicle"]["YEAR"] == "2027"
    assert draft.vehicle_info["NewVehicle"]["MAKE"] == "Ford"
    assert draft.vehicle_info["NewVehicle"]["MODEL"] == "Police Interceptor Utility"


def test_regeneration_refreshes_old_draft_to_project_build_year(tmp_path, monkeypatch):
    paths = _paths(tmp_path)
    draft = new_draft(
        vehicle_info={
            "VehicleType": "PIU",
            "Agency": "Granite Falls Police Department",
            "BuildYear": "2026",
            "ProjectID": "granite-2027",
            "NewVehicle": {
                "MODEL": "2026 Ford Police Interceptor Utility",
                "UNIT ID": "02",
                "YEAR": "2026",
            },
            "ExistingVehicle": {},
        },
        parts=[DraftPart(name="Legacy part", include=False)],
    )
    save_draft(draft, paths.workspace_drafts_dir)
    project = _granite_project(draft.draft_id)
    save_project(project, paths)

    captured: dict = {}

    monkeypatch.setattr("dtm_buildsheet.config.loader.load_configs", lambda paths: object())

    def fake_build_plan(project_input, config):
        captured["info"] = project_input.info
        return BuildPlan(version="test", project=project_input.info, planned_parts=[])

    def fake_render_plan_to_ppt(plan, active_paths):
        output = active_paths.workspace_output_dir / "granite.pptx"
        output.parent.mkdir(parents=True, exist_ok=True)
        output.write_bytes(b"pptx")
        return output

    monkeypatch.setattr("dtm_buildsheet.planning.planner.build_plan", fake_build_plan)
    monkeypatch.setattr("dtm_buildsheet.render_ppt.render_plan_to_ppt", fake_render_plan_to_ppt)
    monkeypatch.setattr(
        "dtm_buildsheet.app.services.exports_upload_service.upload_export_in_background",
        lambda *args, **kwargs: None,
    )

    result = handle_generate_from_draft(
        {"draft_id": draft.draft_id, "project_id": project.project_id}, paths
    )

    assert result["ok"] is True, result.get("error")
    assert captured["info"]["BuildYear"] == "2027"
    assert captured["info"]["NewVehicle"]["YEAR"] == "2027"
    assert captured["info"]["NewVehicle"]["MAKE"] == "Ford"
    assert captured["info"]["NewVehicle"]["MODEL"] == "Police Interceptor Utility"


def test_output_filename_prefers_build_year_over_vehicle_year():
    name = build_output_filename(
        {
            "Agency": "Granite Falls PD",
            "BuildType": "Patrol",
            "BuildYear": "2027",
            "NewVehicle": {"UNIT ID": "02", "YEAR": "2026"},
        }
    )

    assert "2027_GFP_Vehicle_Patrol_Unit_02_" in name
    assert "_2026_Vehicle_Patrol_Unit_02_" not in name


def test_rendered_build_package_uses_build_year(tmp_path, app_paths):
    paths = replace(app_paths, workspace_output_dir=tmp_path / "output")
    paths.workspace_output_dir.mkdir(parents=True)
    plan = BuildPlan(
        version="test",
        project={
            "Agency": "Granite Falls Police Department",
            "BuildType": "Patrol",
            "BuildYear": "2027",
            "VehicleType": "PIU",
            "NewVehicle": {
                "MODEL": "2026 Ford Police Interceptor Utility",
                "UNIT ID": "02",
                "YEAR": "2026",
            },
            "ExistingVehicle": {},
        },
        planned_parts=[],
    )

    output = render_plan_to_ppt(plan, paths)
    assert "2027_GFPD_PIU_Patrol_Unit_02_" in output.name

    from pptx import Presentation

    text = "\n".join(
        shape.text
        for slide in Presentation(str(output)).slides
        for shape in slide.shapes
        if hasattr(shape, "text")
    )
    assert "2027" in text


def test_quickbooks_job_name_prefers_project_build_year():
    project = SimpleNamespace(
        customer=SimpleNamespace(build_year="2027", quote_number="Q-1")
    )
    build_unit = SimpleNamespace(vehicle_model="PIU")
    individual = SimpleNamespace(
        year="2026", model="Police Interceptor Utility", unit_number="02", individual_id="ind-1"
    )

    assert _job_display_name(project, build_unit, individual).startswith("2027 ")
