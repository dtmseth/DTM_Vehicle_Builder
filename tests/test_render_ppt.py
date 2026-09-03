"""Baseline tests: render_ppt produces a valid, openable .pptx file."""
from __future__ import annotations

from pathlib import Path
from types import SimpleNamespace

import pytest
from PIL import Image
from pptx import Presentation
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_SHAPE_TYPE
from pptx.oxml.ns import qn
from pptx.util import Inches

from dtm_buildsheet.models import PartInput, ProjectInput
from dtm_buildsheet.paths import AppPaths
from dtm_buildsheet.planner import build_plan
from dtm_buildsheet.render_ppt import (
    _build_accessory_map,
    _placement_key,
    _render_failures_for_view,
    _stack_parts_above_vehicle,
    _vertical_mirror_slot_is_reflected,
    render_plan_to_ppt,
)
from dtm_buildsheet.ppt_helpers import (
    FOOTER_H,
    SLIDE_H_EMU,
    SLIDE_W_EMU,
    _ManifestEntry,
    _manifest_groups,
    _manifest_page_rows,
    _manifest_row_height,
    _manifest_source_label,
    _badge_label,
    _legend_callouts,
    add_render_exception_slides,
    fill_overview,
    fill_notes,
    MANIFEST_HDR_ROW_H,
    MANIFEST_SECTION_ROW_H,
    physical_light_head_count,
    place_legend,
)


def test_overview_separates_current_and_existing_vehicle_fields():
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    project = SimpleNamespace(
        info={
            "Agency": "Test PD",
            "BuildYear": "2027",
            "BuildType": "Patrol",
            "VehicleType": "PIU",
            "NewVehicle": {
                "YEAR": "2027", "MODEL": "PIU", "UNIT ID": "12",
                "VIN": "ACTUAL123456",
            },
            "ExistingVehicle": {
                "YEAR": "2020", "MAKE": "Ford", "MODEL": "PIU",
                "BUILD TYPE": "Admin", "UNIT ID": "OLD-12",
                "VIN": "OLDVIN654321",
            },
            "UnitNotes": "Mount the radio control head above the siren controller.",
        },
        parts=[],
    )

    fill_overview(slide, project)

    text = "\n".join(
        shape.text for shape in slide.shapes if getattr(shape, "has_text_frame", False)
    )
    assert "ACTUAL123456" in text
    assert "OLDVIN654321" in text
    assert "OLD-12" in text
    assert "Build Type:  Patrol" in text
    assert "Unit #:  12" in text
    assert "Build Type:  Admin" in text
    assert "Unit #:  OLD-12" in text
    assert "UNIT NOTES" in text
    assert "Mount the radio control head above the siren controller." in text


def test_physical_light_head_count_includes_all_supported_input_routes():
    def planned(
        *, name, quantity, render_kind="light", part_id="", components=None,
        placement_counts=(),
    ):
        return SimpleNamespace(
            part_name=name,
            part_id=part_id,
            render_kind=render_kind,
            raw=SimpleNamespace(quantity=quantity, components=components or []),
            placements=[
                SimpleNamespace(instances=[object()] * count)
                for count in placement_counts
            ],
        )

    parts = [
        planned(
            name="Front Warning", quantity=2,
            components=[
                {"part_number": "XI2D", "quantity": 1},
                {"part_number": "XI2E", "quantity": 1},
            ],
            placement_counts=(2, 2),
        ),
        planned(name="Included bumper top-tube lights", quantity=4, placement_counts=(4, 4)),
        planned(name="Included bumper side lights", quantity=2, placement_counts=(2, 2)),
        planned(
            name="Side Warning", quantity=2, part_id="tracer_5lamp",
            placement_counts=(5,),
        ),
        planned(name="Interior Light Bar", quantity=1, placement_counts=(10,)),
        planned(name="Roof Light Bar", quantity=1, render_kind="bar", placement_counts=(1,)),
    ]

    assert physical_light_head_count(parts) == 28

    legacy_fixture = planned(
        name="Side Warning", quantity=5, part_id="tracer_5lamp",
        placement_counts=(5,),
    )
    assert physical_light_head_count([legacy_fixture]) == 5


@pytest.fixture(scope="module")
def rendered_pptx(tmp_path_factory, stearns_input, config):
    out_dir = tmp_path_factory.mktemp("output")
    paths = AppPaths(
        project_root=config.paths.project_root,
        package_dir=config.paths.package_dir,
        resources_dir=config.paths.resources_dir,
        assets_dir=config.paths.assets_dir,
        templates_dir=config.paths.templates_dir,
        workspace_dir=config.paths.workspace_dir,
        workspace_config_dir=config.paths.workspace_config_dir,
        workspace_assets_dir=config.paths.workspace_assets_dir,
        workspace_input_dir=config.paths.workspace_input_dir,
        workspace_output_dir=out_dir,
        samples_dir=config.paths.samples_dir,
    )
    plan = build_plan(stearns_input, config)
    ppt_path = render_plan_to_ppt(plan, paths)
    return ppt_path


def test_render_produces_file(rendered_pptx):
    assert rendered_pptx.exists(), f"Expected output file at {rendered_pptx}"


def test_render_file_has_reasonable_size(rendered_pptx):
    size = rendered_pptx.stat().st_size
    assert size > 10_000, f"Output file is suspiciously small: {size} bytes"


def test_render_file_has_pptx_extension(rendered_pptx):
    assert rendered_pptx.suffix == ".pptx"


def test_render_output_is_openable_by_pptx(rendered_pptx):
    prs = Presentation(str(rendered_pptx))
    assert prs is not None


def test_render_has_slides(rendered_pptx):
    prs = Presentation(str(rendered_pptx))
    assert len(prs.slides) > 0


def test_render_has_multiple_slides(rendered_pptx):
    prs = Presentation(str(rendered_pptx))
    # Expect at least a cover + one vehicle view + manifest
    assert len(prs.slides) >= 3


def test_export_page_order_is_cover_then_vehicle_views_then_manifest_then_notes(rendered_pptx):
    prs = Presentation(str(rendered_pptx))
    slide_text = [
        "\n".join(
            getattr(shape, "text", "")
            for shape in slide.shapes
            if getattr(shape, "has_text_frame", False)
        )
        for slide in prs.slides
    ]
    view_indices = [
        index for index, text in enumerate(slide_text)
        if any(f"{view} VIEW" in text for view in ("FRONT", "SIDE", "TOP", "REAR"))
    ]
    manifest_indices = [
        index for index, text in enumerate(slide_text) if "PARTS MANIFEST" in text
    ]
    notes_index = next(index for index, text in enumerate(slide_text) if "BUILD NOTES" in text)

    assert view_indices
    assert manifest_indices
    assert min(view_indices) == 1
    assert max(view_indices) < min(manifest_indices) < notes_index


def test_manifest_data_text_is_at_least_nine_points(rendered_pptx):
    prs = Presentation(str(rendered_pptx))
    font_sizes = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if not getattr(shape, "has_table", False):
                continue
            table = shape.table
            if not table.rows or table.cell(0, 0).text != "PART / SKU":
                continue
            for row in list(table.rows)[1:]:
                for cell in row.cells:
                    for paragraph in cell.text_frame.paragraphs:
                        for run in paragraph.runs:
                            if run.text.strip():
                                font_sizes.append(run.font.size.pt)

    assert font_sizes
    assert min(font_sizes) >= 9


def test_rendered_manifest_flows_categories_and_never_orphans_section_header(rendered_pptx):
    prs = Presentation(str(rendered_pptx))
    section_names = {
        "LIGHTING", "STRUCTURAL EQUIPMENT", "EQUIPMENT & ELECTRONICS", "OTHER / CUSTOM",
    }
    section_counts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if not getattr(shape, "has_table", False):
                continue
            table = shape.table
            if table.cell(0, 0).text != "PART / SKU":
                continue
            row_labels = [row.cells[0].text.strip() for row in table.rows]
            sections = [
                label for label in row_labels
                if label.removesuffix(" - CONTINUED") in section_names
            ]
            section_counts.append(len(sections))
            assert row_labels[-1].removesuffix(" - CONTINUED") not in section_names

    assert section_counts
    assert any(count >= 2 for count in section_counts)


def test_reference_pages_follow_vehicle_views_before_manifest(tmp_path, stearns_input, config):
    reference_path = tmp_path / "reference.jpg"
    Image.new("RGB", (1200, 800), (35, 75, 115)).save(reference_path, "JPEG")
    output_dir = tmp_path / "output"
    output_dir.mkdir()
    paths = AppPaths(
        project_root=config.paths.project_root,
        package_dir=config.paths.package_dir,
        resources_dir=config.paths.resources_dir,
        assets_dir=config.paths.assets_dir,
        templates_dir=config.paths.templates_dir,
        workspace_dir=config.paths.workspace_dir,
        workspace_config_dir=config.paths.workspace_config_dir,
        workspace_assets_dir=config.paths.workspace_assets_dir,
        workspace_input_dir=config.paths.workspace_input_dir,
        workspace_output_dir=output_dir,
        samples_dir=config.paths.samples_dir,
    )
    plan = build_plan(stearns_input, config)
    plan.reference_photos = [{
        "reference_id": "reference-1",
        "file_name": reference_path.name,
        "title": "Reference detail",
        "note": "Match the installation position shown in this photo.",
        "local_path": str(reference_path),
        "source_relative_path": f"Build Reference Photos/{reference_path.name}",
    }]

    prs = Presentation(str(render_plan_to_ppt(plan, paths)))
    slide_text = [
        "\n".join(
            getattr(shape, "text", "")
            for shape in slide.shapes
            if getattr(shape, "has_text_frame", False)
        )
        for slide in prs.slides
    ]
    view_indices = [
        index for index, text in enumerate(slide_text)
        if any(f"{view} VIEW" in text for view in ("FRONT", "SIDE", "TOP", "REAR"))
    ]
    reference_indices = [
        index for index, text in enumerate(slide_text) if "BUILD REFERENCE PHOTOS" in text
    ]
    manifest_indices = [
        index for index, text in enumerate(slide_text) if "PARTS MANIFEST" in text
    ]
    notes_index = next(index for index, text in enumerate(slide_text) if "BUILD NOTES" in text)

    assert view_indices
    assert reference_indices
    assert manifest_indices
    assert max(view_indices) < min(reference_indices)
    assert max(reference_indices) < min(manifest_indices) < notes_index


def test_cover_and_customer_output_hide_removed_and_internal_fields(rendered_pptx):
    prs = Presentation(str(rendered_pptx))
    cover_text = "\n".join(
        getattr(shape, "text", "")
        for shape in prs.slides[0].shapes
        if getattr(shape, "has_text_frame", False)
    )
    all_text = "\n".join(
        getattr(shape, "text", "")
        for slide in prs.slides
        for shape in slide.shapes
        if getattr(shape, "has_text_frame", False)
    )

    assert "INSTALL TYPE" not in cover_text.upper()
    assert "OTHER ORDERS" not in cover_text.upper()
    assert "QB IMPORT" not in all_text.upper()


def test_build_notes_keep_only_requested_sections_and_do_not_overlap(config):
    prs = Presentation(str(config.paths.templates_dir / "build_sheet_template.pptx"))
    notes_slide = prs.slides[5]
    long_installation_note = (
        "Route the controller harness behind the center console, protect every pass-through, "
        "leave a service loop near the cage, label both ends, and confirm the final routing "
        "with the shop lead before reinstalling the trim panels."
    )
    fill_notes(notes_slide, {
        "PROJECT-WIDE NOTES": ["Keep the vehicle available for inspection."],
        "INSTALLATION NOTES": [long_installation_note],
        "DELIVERY REQUIREMENTS": ["Call fleet before delivery."],
        "CUSTOMER REQUESTS": ["This retired section must not render."],
        "SPECIAL FABRICATION NOTES": ["This retired section must not render."],
        "FINAL APPROVALS": ["This retired section must not render."],
    })

    text_shapes = [
        shape for shape in notes_slide.shapes
        if getattr(shape, "has_text_frame", False)
    ]
    rendered_text = "\n".join(shape.text for shape in text_shapes)
    assert "PROJECT-WIDE NOTES" in rendered_text
    assert "INSTALLATION NOTES" in rendered_text
    assert "DELIVERY REQUIREMENTS" in rendered_text
    assert "CUSTOMER REQUESTS" not in rendered_text
    assert "SPECIAL FABRICATION" not in rendered_text
    assert "FINAL APPROVALS" not in rendered_text

    install_note_shape = next(
        shape for shape in text_shapes if "Route the controller harness" in shape.text
    )
    delivery_header = next(
        shape for shape in text_shapes if shape.text.strip() == "DELIVERY REQUIREMENTS"
    )
    assert install_note_shape.top + install_note_shape.height <= delivery_header.top


@pytest.mark.parametrize(
    ("anchor_y", "expected"),
    [
        (0.17, [False, True]),
        (0.83, [True, False]),
        (0.50, [False, True]),
    ],
)
def test_vertical_mirror_reflects_only_the_non_anchor_slot(anchor_y, expected):
    placement = SimpleNamespace(
        pattern="vertical_mirror",
        anchor={"y": anchor_y},
        instances=[SimpleNamespace(), SimpleNamespace()],
    )

    assert [
        _vertical_mirror_slot_is_reflected(placement, index)
        for index in range(len(placement.instances))
    ] == expected


def test_render_second_sample(tmp_path_factory, test_build_input, config):
    out_dir = tmp_path_factory.mktemp("output2")
    paths = AppPaths(
        project_root=config.paths.project_root,
        package_dir=config.paths.package_dir,
        resources_dir=config.paths.resources_dir,
        assets_dir=config.paths.assets_dir,
        templates_dir=config.paths.templates_dir,
        workspace_dir=config.paths.workspace_dir,
        workspace_config_dir=config.paths.workspace_config_dir,
        workspace_assets_dir=config.paths.workspace_assets_dir,
        workspace_input_dir=config.paths.workspace_input_dir,
        workspace_output_dir=out_dir,
        samples_dir=config.paths.samples_dir,
    )
    plan = build_plan(test_build_input, config)
    ppt_path = render_plan_to_ppt(plan, paths)
    assert ppt_path.exists()
    prs = Presentation(str(ppt_path))
    assert len(prs.slides) > 0


def test_dual_shroud_groups_pairs_without_dark_outline_box(tmp_path, config):
    paths = AppPaths(
        project_root=config.paths.project_root,
        package_dir=config.paths.package_dir,
        resources_dir=config.paths.resources_dir,
        assets_dir=config.paths.assets_dir,
        templates_dir=config.paths.templates_dir,
        workspace_dir=config.paths.workspace_dir,
        workspace_config_dir=config.paths.workspace_config_dir,
        workspace_assets_dir=config.paths.workspace_assets_dir,
        workspace_input_dir=config.paths.workspace_input_dir,
        workspace_output_dir=tmp_path,
        samples_dir=config.paths.samples_dir,
    )
    parent = PartInput(
        name="Rear Warning 1", part_number="TST0D", part_type="warning_light",
        location="LOWER CARGO WINDOW", raw_color="Red/White", quantity=4,
        line_id="t-series-parent",
        picker_config={
            "mode": "uniform", "colorsPerHead": "duo",
            "uniform": ["red", "white"], "count": 4,
        },
    )
    shroud = PartInput(
        name="Rear Warning 1 · T-Series Shroud", part_number="THSG2", quantity=2,
        line_id="dual-shrouds", parent_line_id=parent.line_id,
        accessory_category="shroud",
        picker_config={"accessory_quantity": {
            "parent_units_per_item": 2, "render_parent_group": "dual_shroud",
        }},
    )
    plan = build_plan(ProjectInput(
        info={"VehicleType": "PIU", "ProjectID": "DUAL-SHROUD"},
        parts=[parent, shroud], notes={},
    ), config)

    prs = Presentation(str(render_plan_to_ppt(plan, paths)))
    compound_groups = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.shape_type != MSO_SHAPE_TYPE.GROUP:
                continue
            children = list(shape.shapes)
            has_housing = any(
                child.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE
                and child.auto_shape_type == MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE
                for child in children
            )
            picture_count = sum(child.shape_type == MSO_SHAPE_TYPE.PICTURE for child in children)
            if picture_count >= 2:
                compound_groups.append((shape, has_housing, picture_count))

    assert compound_groups, "expected the paired T-Series placement to remain grouped"
    assert all(not has_housing for _shape, has_housing, _count in compound_groups)
    assert sum(picture_count for _shape, _housing, picture_count in compound_groups) >= 4


def test_concealed_speaker_renders_translucent_with_mount_callout(tmp_path, config):
    paths = AppPaths(
        project_root=config.paths.project_root,
        package_dir=config.paths.package_dir,
        resources_dir=config.paths.resources_dir,
        assets_dir=config.paths.assets_dir,
        templates_dir=config.paths.templates_dir,
        workspace_dir=config.paths.workspace_dir,
        workspace_config_dir=config.paths.workspace_config_dir,
        workspace_assets_dir=config.paths.workspace_assets_dir,
        workspace_input_dir=config.paths.workspace_input_dir,
        workspace_output_dir=tmp_path,
        samples_dir=config.paths.samples_dir,
    )
    speaker = PartInput(
        name="Siren Speaker", part_number="SA315P", part_type="siren_speaker",
        location="BEHIND GRILL (CENTER)", quantity=1, line_id="concealed-speaker",
    )
    plan = build_plan(ProjectInput(
        info={"VehicleType": "PIU", "ProjectID": "CONCEALED-SPEAKER"},
        parts=[speaker], notes={},
    ), config)

    prs = Presentation(str(render_plan_to_ppt(plan, paths)))
    front_slide = next(
        slide for slide in prs.slides
        if any("FRONT VIEW" in getattr(shape, "text", "") for shape in slide.shapes)
    )
    alpha_values = [
        alpha.get("amt")
        for shape in front_slide.shapes
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE
        for alpha in shape._element.findall(".//" + qn("a:alphaModFix"))
    ]

    assert "80000" in alpha_values
    callout = next(
        shape for shape in front_slide.shapes
        if getattr(shape, "text", "").strip() == "SPEAKER BEHIND GRILLE"
    )
    fill_alpha = callout._element.find(".//" + qn("a:solidFill") + "//" + qn("a:alpha"))
    assert fill_alpha is not None
    assert fill_alpha.get("val") == "70000"
    assert callout.width < 1.42 * 914400
    assert callout.height == int(0.18 * 914400)
    leaders = [
        shape for shape in front_slide.shapes
        if shape.shape_type == MSO_SHAPE_TYPE.LINE
    ]
    assert leaders
    assert any(leader.line.color.rgb == callout.fill.fore_color.rgb for leader in leaders)


def test_concealed_speaker_callout_manual_offset_changes_export_position(tmp_path, config):
    paths = AppPaths(
        project_root=config.paths.project_root,
        package_dir=config.paths.package_dir,
        resources_dir=config.paths.resources_dir,
        assets_dir=config.paths.assets_dir,
        templates_dir=config.paths.templates_dir,
        workspace_dir=config.paths.workspace_dir,
        workspace_config_dir=config.paths.workspace_config_dir,
        workspace_assets_dir=config.paths.workspace_assets_dir,
        workspace_input_dir=config.paths.workspace_input_dir,
        workspace_output_dir=tmp_path,
        samples_dir=config.paths.samples_dir,
    )
    speaker = PartInput(
        name="Siren Speaker", part_number="SA315P", part_type="siren_speaker",
        location="BEHIND GRILL (CENTER)", quantity=1, line_id="concealed-speaker",
    )
    plan = build_plan(ProjectInput(
        info={"VehicleType": "PIU", "ProjectID": "CONCEALED-SPEAKER-OFFSET"},
        parts=[speaker], notes={},
    ), config)
    placement = plan.planned_parts[0].placements[0]
    placement.callout_dx = 0.10
    placement.callout_dy = -0.05

    prs = Presentation(str(render_plan_to_ppt(plan, paths)))
    front_slide = next(
        slide for slide in prs.slides
        if any("FRONT VIEW" in getattr(shape, "text", "") for shape in slide.shapes)
    )
    callout = next(
        shape for shape in front_slide.shapes
        if getattr(shape, "text", "").strip() == "SPEAKER BEHIND GRILLE"
    )
    speaker_picture = next(
        shape for shape in front_slide.shapes
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE
        and shape._element.find(".//" + qn("a:alphaModFix")) is not None
    )

    assert callout.left + callout.width // 2 > speaker_picture.left + speaker_picture.width // 2
    assert any(shape.shape_type == MSO_SHAPE_TYPE.LINE for shape in front_slide.shapes)


def test_custom_location_points_share_one_legend_card(tmp_path, config):
    """Edina regression: six visual heads are one qty-six legend entry."""
    paths = AppPaths(
        project_root=config.paths.project_root,
        package_dir=config.paths.package_dir,
        resources_dir=config.paths.resources_dir,
        assets_dir=config.paths.assets_dir,
        templates_dir=config.paths.templates_dir,
        workspace_dir=config.paths.workspace_dir,
        workspace_config_dir=config.paths.workspace_config_dir,
        workspace_assets_dir=config.paths.workspace_assets_dir,
        workspace_input_dir=config.paths.workspace_input_dir,
        workspace_output_dir=tmp_path,
        samples_dir=config.paths.samples_dir,
    )
    points = [
        {"x": 0.345 + index * 0.06, "y": 0.621, "head_index": index}
        for index in range(6)
    ]
    part = PartInput(
        name="Warning Light 1", part_number="ION", part_type="warning_light",
        raw_color="Red/White / Blue/White", driver_color="Red/White",
        passenger_color="Blue/White", location="Front Grill", quantity=6,
        line_id="edina-front-grille",
        components=[
            {"part_number": "XI2D", "color": "Red/White", "quantity": 3},
            {"part_number": "XI2E", "color": "Blue/White", "quantity": 3},
        ],
        picker_config={"mode": "split", "custom_location": {
            "label": "Front Grill", "placements": {"front": points},
        }},
    )
    plan = build_plan(ProjectInput(
        info={"VehicleType": "PIU", "ProjectID": "EDINA-GRILLE-LEGEND"},
        parts=[part], notes={},
    ), config)
    assert sum(
        len(placement.instances)
        for placement in plan.planned_parts[0].placements
        if placement.view == "front"
    ) == 6

    prs = Presentation(str(render_plan_to_ppt(plan, paths)))
    front_slide = next(
        slide for slide in prs.slides
        if any("FRONT VIEW" in getattr(shape, "text", "") for shape in slide.shapes)
    )
    legend_cards = [
        shape.text for shape in front_slide.shapes
        if "Front Grill" in getattr(shape, "text", "") and "QTY:" in shape.text
    ]
    assert len(legend_cards) == 1
    assert "QTY: 6" in legend_cards[0]


def test_render_failure_legend_only_lists_failed_current_view_components():
    """Diagram callouts are render diagnostics, not a second manifest."""
    def part(name, *, location, warnings=(), placements=(), line_id=""):
        return SimpleNamespace(
            part_name=name,
            raw=SimpleNamespace(
                line_id=line_id,
                location=location,
                notes="",
                part_number="",
            ),
            warnings=list(warnings),
            placements=list(placements),
        )

    asset_failure = SimpleNamespace(
        view="front", location_key="BUMPER", color_profile="standard", line_id="asset-line"
    )
    top_failure = SimpleNamespace(
        view="top", location_key="ROOF", color_profile="standard", line_id="top-line"
    )
    asset_part = part(
        "Front Scene", location="BUMPER", placements=[asset_failure], line_id="asset-line"
    )
    plan = SimpleNamespace(planned_parts=[
        part("Manifest-only item", location="EQUIPMENT TRAY", warnings=["No views configured"]),
        part("Missing fixture", location="WINDSHIELD", warnings=["front: no fixture coords"]),
        asset_part,
        part("Top-only failure", location="ROOF", placements=[top_failure], line_id="top-line"),
    ])

    items = _render_failures_for_view(
        plan,
        "front",
        {_placement_key(asset_part, asset_failure): ["Missing asset (front): assets/missing.png"]},
    )

    assert [(item.name, item.location) for item in items] == [
        ("Missing fixture", "WINDSHIELD"),
        ("Front Scene", "BUMPER"),
    ]
    assert items[0].notes == "front: no fixture coords"
    assert items[1].notes == "Missing asset (front): assets/missing.png"


def test_manifest_groups_use_separate_customer_facing_categories():
    def planned(name, *, category="", render_kind="equipment", part_type=""):
        return SimpleNamespace(
            category=category,
            render_kind=render_kind,
            raw=SimpleNamespace(
                name=name,
                include=True,
                notes="",
                part_number="",
                location="",
                part_type=part_type,
            ),
            placements=[],
        )

    groups = _manifest_groups([
        planned("Front Warning", category="warning", render_kind="light"),
        planned("Pit Bar", part_type="pit_bar"),
        planned("Push Bumper", part_type="push_bumper"),
        planned("Radio Head", part_type="radio_head"),
        planned("Shop-Supplied Item", category="custom", render_kind="none", part_type="custom_part"),
    ])

    assert [(label, [entry.name for entry in entries]) for label, entries in groups] == [
        ("Lighting", ["Front Warning"]),
        ("Structural Equipment", ["Push Bumper", "Pit Bar"]),
        ("Equipment & Electronics", ["Radio Head"]),
        ("Other / Custom", ["Shop-Supplied Item"]),
    ]


def test_manifest_categories_flow_on_one_page_with_clear_section_rows():
    raw = SimpleNamespace(
        new_or_used="New", source="", supply_type="new",
        customer_condition="", customer_source="",
    )
    groups = [
        ("Lighting", [_ManifestEntry(raw=raw, location="Grille", name="Warning light")]),
        ("Structural Equipment", [_ManifestEntry(raw=raw, location="Front", name="Push bumper")]),
        ("Equipment & Electronics", [_ManifestEntry(raw=raw, location="Console", name="Control head")]),
    ]

    pages = _manifest_page_rows(groups, Inches(6.0))

    assert len(pages) == 1
    assert [value for kind, value, _height in pages[0] if kind == "section"] == [
        "Lighting", "Structural Equipment", "Equipment & Electronics",
    ]
    assert pages[0][-1][0] == "entry"


def test_manifest_section_header_moves_with_first_item_to_next_page():
    raw = SimpleNamespace(
        new_or_used="New", source="", supply_type="new",
        customer_condition="", customer_source="",
    )
    first = _ManifestEntry(raw=raw, location="Grille", name="Warning light")
    second = _ManifestEntry(raw=raw, location="Front", name="Push bumper")
    first_height = _manifest_row_height(first)
    second_height = _manifest_row_height(second)
    just_too_short = (
        MANIFEST_HDR_ROW_H
        + MANIFEST_SECTION_ROW_H + first_height
        + MANIFEST_SECTION_ROW_H + second_height
        - 1
    )

    pages = _manifest_page_rows([
        ("Lighting", [first]),
        ("Structural Equipment", [second]),
    ], just_too_short)

    assert len(pages) == 2
    assert all(page[-1][0] == "entry" for page in pages)
    assert pages[1][0][0:2] == ("section", "Structural Equipment")


def test_manifest_groups_order_lighting_by_function_before_name():
    def planned(name, part_type):
        return SimpleNamespace(
            category="warning",
            render_kind="light",
            raw=SimpleNamespace(
                name=name,
                include=True,
                notes="",
                part_number="",
                location="",
                part_type=part_type,
            ),
            placements=[],
        )

    groups = _manifest_groups([
        planned("Zulu Tracer", "tracer_2_lamp"),
        planned("Alpha Interior", "front_interior_light_bar"),
        planned("Bravo Warning", "warning_light"),
        planned("Charlie Scene", "front_scene"),
        planned("Delta Bar", "roof_light_bar"),
    ])

    assert [entry.name for entry in groups[0][1]] == [
        "Bravo Warning", "Charlie Scene", "Alpha Interior", "Delta Bar", "Zulu Tracer",
    ]


def test_manifest_promotes_selected_skus_and_keeps_guided_components_and_accessories():
    def planned(name, *, line_id, parent_line_id="", components=(), part_number="", part_type="warning_light",
                location="GRILL", lens="", comment=""):
        return SimpleNamespace(
            category="warning",
            render_kind="light",
            placements=[],
            raw=SimpleNamespace(
                name=name,
                include=True,
                notes="Install with supplied hardware",
                comment=comment,
                part_number=part_number,
                location=location,
                part_type=part_type,
                line_id=line_id,
                parent_line_id=parent_line_id,
                components=list(components),
                picker_config={},
                manufacturer="Whelen",
                raw_color="",
                lens=lens,
                quantity=1,
                new_or_used="New",
                source="",
            ),
        )

    parent = planned(
        "Forward Warning", line_id="parent", part_number="ION",
        components=[
            {"part_number": "IONDUO", "quantity": 2, "color": "Red/Blue"},
            {"label": "Mounting location", "location": "Upper grill", "detail": "Centered"},
        ],
        lens="Smoked", comment="Confirm final aiming with the customer.",
    )
    accessory = planned(
        "Forward Warning · Bracket", line_id="child", parent_line_id="parent",
        part_number="BRKT-1", part_type="bracket_mount",
    )

    groups = _manifest_groups([accessory, parent])

    assert [(entry.name, entry.part_number, entry.location, entry.indent) for entry in groups[0][1]] == [
        ("Forward Warning", "IONDUO", "GRILL", 0),
        ("Mounting location", "", "Upper grill", 1),
        ("Bracket", "BRKT-1", "GRILL", 1),
    ]
    selected_sku = groups[0][1][0]
    assert "Lens: Smoked" in selected_sku.detail
    assert selected_sku.comment == "Confirm final aiming with the customer."


def test_manifest_keeps_exact_console_base_and_included_motion_component_visible():
    def planned(name, *, line_id, parent_line_id="", part_number="", part_type="console", picker_config=None):
        return SimpleNamespace(
            category="structural",
            render_kind="none",
            placements=[],
            raw=SimpleNamespace(
                name=name,
                include=True,
                notes="",
                comment="",
                part_number=part_number,
                location="IN CENTER CONSOLE",
                part_type=part_type,
                line_id=line_id,
                parent_line_id=parent_line_id,
                components=[],
                picker_config=picker_config or {},
                manufacturer="Gamber Johnson",
                raw_color="",
                lens="",
                quantity=1,
                new_or_used="New",
                source="",
            ),
        )

    parent = planned(
        "Center Console", line_id="console", part_number="7170-0734-04",
    )
    motion = planned(
        "Center Console · Motion Attachment", line_id="motion", parent_line_id="console",
        part_number="7160-0220", part_type="motion_attachment",
        picker_config={"console_kit_included": True},
    )

    entries = _manifest_groups([motion, parent])[0][1]

    assert [(entry.name, entry.part_number, entry.indent) for entry in entries] == [
        ("Center Console", "7170-0734-04", 0),
        ("Motion Attachment", "7160-0220", 1),
    ]


def test_manifest_combines_standard_duo_skus_into_one_compact_shop_row():
    raw = SimpleNamespace(
        name="Forward Warning 1", include=True, notes="", comment="Aim evenly",
        part_number="ION", location="UPPER GRILLE", part_type="warning_light",
        line_id="duo", parent_line_id="", manufacturer="Whelen",
        raw_color="Red/White / Blue/White", driver_color="Red/White",
        passenger_color="Blue/White", lens="Smoked", quantity=4,
        new_or_used="New", source="",
        picker_config={"colorsPerHead": "duo", "mode": "split"},
        components=[
            {"part_number": "3SBCCDCR", "quantity": 2, "color": "Blue/White"},
            {"part_number": "3SRCCDCR", "quantity": 2, "color": "Red/White"},
        ],
    )
    planned = SimpleNamespace(
        category="warning", render_kind="light", placements=[], raw=raw,
    )

    entries = _manifest_groups([planned])[0][1]

    assert len(entries) == 1
    assert entries[0].name == "Forward Warning 1"
    assert entries[0].part_number == "3SBCCDCR / 3SRCCDCR"
    assert entries[0].quantity == 4
    assert entries[0].location == "UPPER GRILLE"
    assert entries[0].detail == "Blue/White (passenger)  ·  Red/White (driver)  ·  Lens: Smoked"
    assert "\n" not in entries[0].detail


def test_manifest_uses_live_round_light_quantity_without_stale_qb_recipe(config):
    raw = SimpleNamespace(
        name="Rear Seat Cargo Lights", include=True, notes="", comment="",
        part_number="3SRCCDCR", location="LOWER KICK PANELS",
        part_type="rear_seat_cargo_lights", line_id="round-live",
        parent_line_id="", manufacturer="Unassigned (QB Import)",
        raw_color="Red/White", lens="Clear", quantity=2,
        new_or_used="New", source="", components=[],
        picker_config={},
    )
    planned = SimpleNamespace(
        category="warning", render_kind="light", placements=[], raw=raw,
    )

    entry = _manifest_groups([planned], config.paths)[0][1][0]

    assert entry.quantity == 2
    assert entry.description == 'WHELEN 3" ROUND SPLIT RED/WHT COMPART'
    assert "x2 REAR CARGO" not in entry.description
    assert "x3 CAGE" not in entry.description
    assert entry.manufacturer == "Unassigned"


def test_manifest_window_tint_keeps_percentage_but_hides_pricing(config):
    raw = SimpleNamespace(
        name="Window Tint", include=True, notes="", comment="",
        part_number="TINT", location="SELECTED WINDOWS", part_type="window_tint",
        line_id="tint", parent_line_id="", manufacturer="Unassigned (QB Import)",
        raw_color="", lens="", quantity=3, new_or_used="New", source="",
        components=[], picker_config={
            "window_tint": {
                "windows": ["windshield_brow", "driver_front", "passenger_front"],
                "percentage": 35,
                "unit_price": 65,
            },
        },
    )
    planned = SimpleNamespace(
        category="custom", render_kind="none", placements=[], raw=raw,
    )

    entry = _manifest_groups([planned], config.paths)[0][1][0]

    assert "35% tint" in entry.detail
    assert "$" not in entry.detail
    assert "65" not in entry.detail
    assert entry.manufacturer == "Unassigned"


def test_build_legend_groups_direct_child_parts_with_their_parent():
    def planned(name, *, line_id, parent_line_id="", part_number="", accessory_of=None):
        return SimpleNamespace(
            part_name=name,
            accessory_of=accessory_of,
            raw=SimpleNamespace(
                line_id=line_id,
                parent_line_id=parent_line_id,
                part_number=part_number,
                notes="",
            ),
        )

    parent = planned("Control Head 1", line_id="control-head", part_number="CCTL7")
    child = planned(
        "Control Head 1 · Magnetic Mic", line_id="mic", parent_line_id="control-head",
        part_number="MMSU-1",
    )
    legacy_child = planned(
        "Siren Amplifier Harness", line_id="legacy", part_number="SAH-1",
        accessory_of="Siren Amplifier",
    )

    assert _build_accessory_map(SimpleNamespace(
        planned_parts=[parent, child, legacy_child]
    )) == {
        "control-head": [("Magnetic Mic", "MMSU-1")],
        "Siren Amplifier": [("Siren Amplifier Harness", "SAH-1")],
    }


def test_build_legend_card_displays_children_by_parent_line_id():
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    parent = SimpleNamespace(
        name="Control Head 1", line_id="control-head", location="",
        manufacturer="Whelen", part_number="CCTL7", color="", lens="",
        new_or_used="New", source="", is_reused=False, quantity=2,
    )

    place_legend(
        slide, [parent], [],
        {"control-head": [("Magnetic Mic", "MMSU-1")]},
        view="front",
    )

    text = "\n".join(
        shape.text for shape in slide.shapes if getattr(shape, "has_text_frame", False)
    )
    assert "+ Magnetic Mic  ·  MMSU-1" in text
    assert "QTY: 2" in text
    card = next(
        shape for shape in slide.shapes
        if "+ Magnetic Mic" in getattr(shape, "text", "")
    )
    supporting_runs = [
        run for paragraph in card.text_frame.paragraphs[1:]
        for run in paragraph.runs if run.text.strip()
    ]
    assert supporting_runs
    assert all(str(run.font.color.rgb) == "000000" for run in supporting_runs)


def test_customer_supplied_manifest_source_is_compact_and_explicit():
    raw = SimpleNamespace(
        supply_type="customer_supplied", customer_condition="used",
        customer_source="", new_or_used="Used", source="",
    )
    entry = _ManifestEntry(
        raw=raw, location="CONSOLE", name="Customer radio", quantity=1,
    )

    assert _manifest_source_label(raw) == "Customer supplied\nUsed - Source needed"
    assert _manifest_row_height(entry) == Inches(0.38)


def test_manifest_hides_fixture_ids_and_locations_repeating_the_part_name():
    def planned(name, *, location, placements=()):
        return SimpleNamespace(
            category="warning",
            render_kind="light",
            placements=list(placements),
            raw=SimpleNamespace(
                name=name, include=True, notes="", comment="", part_number="IONDUO",
                location=location, part_type="warning_light", line_id=name,
                parent_line_id="", components=[], picker_config={}, manufacturer="Whelen",
                raw_color="Red/Blue", lens="Smoked", quantity=1, new_or_used="New", source="",
            ),
        )

    fixture = planned("Tracer", location="FIXTURE:TRACER_5LAMP")
    repeated = planned(
        "Front Interior Light Bar · Inner Edge Lighthead",
        location="Interior Light Bar (Front)",
    )

    entries = _manifest_groups([fixture, repeated])[0][1]

    assert [entry.location for entry in entries] == ["", ""]


def test_manifest_groups_bumper_and_siren_systems_and_hides_picker_process_details():
    def planned(name, *, part_type, notes="", components=(), part_number="", quantity=1):
        return SimpleNamespace(
            category="equipment",
            render_kind="equipment",
            placements=[],
            raw=SimpleNamespace(
                name=name, include=True, notes=notes, comment="", part_number=part_number,
                location="", part_type=part_type, line_id=name, parent_line_id="",
                components=list(components), picker_config={}, manufacturer="Whelen",
                raw_color="", lens="", quantity=quantity, new_or_used="New", source="",
            ),
        )

    groups = _manifest_groups([
        planned("Wing Wraps", part_type="wing_wraps"),
        planned("Pit Bar", part_type="pit_bar"),
        # Legacy rows lack the newer part_type identifier; they still need to
        # sit with their bumper-system peers in the customer-facing manifest.
        planned("Push Bumper", part_type=""),
        planned("Howler", part_type="howler"),
        planned("Siren Speaker 1", part_type="siren_speaker", part_number="SA315P"),
        planned(
            "Control Head 1", part_type="control_head",
            notes="PA mic: Driver's door · Magnetic mic",
            components=[
                {"part_number": "CCTL7"},
                {"label": "PA Mic", "location": "Driver's door", "detail": "Magnetic mic"},
            ],
        ),
        planned(
            "Radio Communications", part_type="radio_head",
            notes="Radio Communications — guided system details",
            components=[
                {"label": "Radio microphone", "detail": "Uses the selected center-console mic clip"},
                {"label": "Radio speaker", "location": "Back of center console", "detail": "Shop mounting location"},
            ],
        ),
    ])
    entries = {label: rows for label, rows in groups}

    structural_names = [entry.name for entry in entries["Structural Equipment"]]
    assert structural_names[:3] == ["Push Bumper", "Pit Bar", "Wing Wraps"]
    equipment_names = [entry.name for entry in entries["Equipment & Electronics"]]
    assert equipment_names[:2] == ["Siren Speaker(s)", "Howler"]

    control_head_rows = [entry for entry in entries["Equipment & Electronics"] if entry.name == "Control Head 1"]
    assert len(control_head_rows) == 1
    assert "PA mic: Driver's door" in control_head_rows[0].detail
    assert "PA Mic" not in equipment_names
    assert "guided system" not in "\n".join(
        f"{entry.description}\n{entry.detail}" for entry in entries["Equipment & Electronics"]
    ).casefold()
    radio_speaker = next(entry for entry in entries["Equipment & Electronics"] if entry.name == "Radio speaker")
    assert radio_speaker.location == "Back of center console"
    assert radio_speaker.detail == ""


def test_ppt_vehicle_stays_below_negative_layer_bumper_and_lights():
    from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    vehicle = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, 1, 1)
    bumper = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, 1, 1)
    light = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, 1, 1)

    _stack_parts_above_vehicle(
        slide, vehicle._element,
        [(light._element, 0), (bumper._element, -20)],
    )

    elements = list(slide.shapes._spTree)
    assert elements.index(vehicle._element) < elements.index(bumper._element) < elements.index(light._element)


def test_blank_legacy_status_maps_to_new_even_when_old_source_text_is_present():
    text, _ = _badge_label(SimpleNamespace(
        new_or_used="", source="Customer supplied", is_reused=False,
    ))

    assert text == "■ NEW"


def test_visual_part_callouts_show_comment_and_used_source():
    part = SimpleNamespace(
        supply_type="customer_supplied", customer_condition="used",
        customer_source="Retired Unit 12", new_or_used="Used", source="Retired Unit 12",
        comment="Confirm antenna routing with customer.",
    )

    assert _legend_callouts(part) == [
        "NOTE: Confirm antenna routing with customer.",
        "USED SOURCE: Retired Unit 12",
    ]


def test_render_exception_detail_pages_paginate_without_crossing_footer():
    prs = Presentation()
    prs.slide_width = SLIDE_W_EMU
    prs.slide_height = SLIDE_H_EMU
    failures = [
        SimpleNamespace(
            name=f"Unresolved Front Component {index}",
            location="FRONT BUMPER / WINDSHIELD",
            notes="front: The placement asset and coordinates could not be resolved for this item.",
        )
        for index in range(30)
    ]

    pages_added = add_render_exception_slides(
        prs, [("front", failures)], paths=None, footer_text="Test footer"
    )

    assert pages_added >= 2
    rendered_rows = 0
    for slide in prs.slides:
        tables = [shape for shape in slide.shapes if shape.has_table]
        assert len(tables) == 1
        table_shape = tables[0]
        assert table_shape.top + table_shape.height <= prs.slide_height - FOOTER_H
        rendered_rows += len(table_shape.table.rows) - 1  # exclude repeated header
    assert rendered_rows == len(failures)
