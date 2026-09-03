from __future__ import annotations

import json
from pathlib import Path

import pytest
from PIL import Image
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.oxml.ns import qn
from pptx.util import Inches

from dtm_buildsheet.app.services.reference_media_service import resolve_reference_media
from dtm_buildsheet.domain.plan_models import BuildPlan
from dtm_buildsheet.domain.project_models import BuildReferenceAsset
from dtm_buildsheet.domain.reference_photo_layout import plan_reference_photo_pages
from dtm_buildsheet.paths import AppPaths
from dtm_buildsheet.ppt_helpers import (
    _safe_relative_photo_link,
    add_reference_photo_slides,
)


def _paths(tmp_path: Path) -> AppPaths:
    assets = tmp_path / "assets"
    assets.mkdir()
    workspace = tmp_path / "workspace"
    workspace.mkdir()
    return AppPaths(
        workspace_dir=workspace,
        workspace_assets_dir=assets,
        assets_dir=assets,
    )


def _photo(path: Path, color: tuple[int, int, int], size=(1200, 800)) -> None:
    Image.new("RGB", size, color).save(path, "JPEG")


@pytest.mark.parametrize(("photo_count", "expected_pages"), [(1, 1), (4, 1), (5, 2)])
def test_reference_photos_render_four_per_page(tmp_path, photo_count, expected_pages):
    paths = _paths(tmp_path)
    photos = []
    for index in range(photo_count):
        source = tmp_path / f"photo-{index + 1}.jpg"
        _photo(source, (30 + index * 10, 80, 120))
        photos.append({
            "reference_id": f"ref-{index + 1}",
            "file_name": source.name,
            "title": f"Photo {index + 1}",
            "note": f"Install detail {index + 1}",
            "local_path": str(source),
            "source_relative_path": f"Build Reference Photos/{source.name}",
        })
    plan = BuildPlan(
        version="test",
        project={
            "Agency": "Test PD",
            "BuildType": "Patrol",
            "CanonicalVehicleName": "2027 PIU - Patrol - Unit 12",
            "NewVehicle": {"UNIT ID": "12"},
        },
        planned_parts=[],
        reference_photos=photos,
    )
    presentation = Presentation()

    added = add_reference_photo_slides(presentation, plan, paths, footer_text="Test footer")

    assert added == expected_pages
    assert len(presentation.slides) == expected_pages
    pictures = [
        shape for slide in presentation.slides for shape in slide.shapes
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE
    ]
    assert len(pictures) == photo_count
    assert pictures[0].click_action.hyperlink.address == "Build Reference Photos/photo-1.jpg"
    text = "\n".join(
        shape.text for slide in presentation.slides for shape in slide.shapes
        if getattr(shape, "has_text_frame", False)
    )
    assert "BUILD REFERENCE PHOTOS" in text
    assert "2027 PIU - Patrol - Unit 12" in text
    assert "Install detail 1" in text


def test_reference_pages_omitted_when_empty(tmp_path):
    presentation = Presentation()
    plan = BuildPlan(version="test", project={}, planned_parts=[])
    assert add_reference_photo_slides(presentation, plan, _paths(tmp_path)) == 0
    assert len(presentation.slides) == 0


def test_reference_photo_uses_large_image_with_readable_overlay(tmp_path):
    paths = _paths(tmp_path)
    source = tmp_path / "overlay.jpg"
    _photo(source, (35, 75, 115))
    plan = BuildPlan(
        version="test",
        project={"CanonicalVehicleName": "2027 PIU - Patrol - Unit 12"},
        planned_parts=[],
        reference_photos=[{
            "reference_id": "overlay",
            "file_name": source.name,
            "title": "Push bumper alignment",
            "note": "Match the bracket position shown here.",
            "local_path": str(source),
            "source_relative_path": f"Build Reference Photos/{source.name}",
        }],
    )
    presentation = Presentation()

    assert add_reference_photo_slides(presentation, plan, paths) == 1
    slide = presentation.slides[0]
    picture = next(shape for shape in slide.shapes if shape.shape_type == MSO_SHAPE_TYPE.PICTURE)
    overlay = next(
        shape for shape in slide.shapes
        if "Match the bracket position" in getattr(shape, "text", "")
    )
    fill_alpha = overlay._element.find(".//" + qn("a:solidFill") + "//" + qn("a:alpha"))
    font_sizes = [
        run.font.size.pt for paragraph in overlay.text_frame.paragraphs
        for run in paragraph.runs if run.text.strip()
    ]

    assert picture.height >= Inches(5.75)
    assert overlay.left == picture.left
    assert overlay.width == picture.width
    assert overlay.top + overlay.height == picture.top + picture.height
    assert fill_alpha is not None and fill_alpha.get("val") == "78000"
    assert font_sizes == [11, 10]
    assert overlay.click_action.hyperlink.address == "Build Reference Photos/overlay.jpg"


def test_adaptive_layout_gives_portrait_a_full_height_column(tmp_path):
    paths = _paths(tmp_path)
    sizes = [(800, 1200), (1200, 800), (1200, 800)]
    photos = []
    for index, size in enumerate(sizes):
        source = tmp_path / f"mixed-{index}.jpg"
        _photo(source, (40 + index * 20, 80, 120), size)
        photos.append({
            "reference_id": f"mixed-{index}",
            "file_name": source.name,
            "local_path": str(source),
            "source_relative_path": f"Build Reference Photos/{source.name}",
        })
    plan = BuildPlan(version="test", project={}, planned_parts=[], reference_photos=photos)
    presentation = Presentation()

    assert add_reference_photo_slides(presentation, plan, paths) == 1
    pictures = [
        shape for shape in presentation.slides[0].shapes
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE
    ]
    assert len(pictures) == 3
    assert pictures[0].height > pictures[1].height * 2
    assert pictures[0].top <= pictures[1].top
    assert pictures[2].top > pictures[1].top


def test_adaptive_page_planning_keeps_three_landscapes_before_portrait():
    pages = plan_reference_photo_pages([
        (1200, 800), (1200, 800), (1200, 800), (800, 1200),
    ])

    assert [len(page.placements) for page in pages] == [3, 1]
    portrait = pages[1].placements[0]
    assert portrait.height == 1
    assert portrait.width < 1


def test_serialized_plan_never_contains_local_reference_path(tmp_path):
    plan = BuildPlan(
        version="test",
        project={},
        planned_parts=[],
        reference_photos=[{
            "reference_id": "ref-1",
            "file_name": "one.jpg",
            "local_path": str(tmp_path / "private" / "one.jpg"),
            "source_relative_path": "Build Reference Photos/one.jpg",
        }],
    )
    serialized = plan.to_dict()
    assert "local_path" not in serialized["reference_photos"][0]
    assert serialized["reference_photos"][0]["source_relative_path"] == "Build Reference Photos/one.jpg"


def test_relative_photo_link_rejects_escape_or_nested_paths():
    assert _safe_relative_photo_link("Build Reference Photos/one.jpg") == "Build Reference Photos/one.jpg"
    assert _safe_relative_photo_link("../one.jpg") == ""
    assert _safe_relative_photo_link("Build Reference Photos/nested/one.jpg") == ""


def test_reference_media_uses_matching_local_cache_cloud_off(tmp_path):
    paths = _paths(tmp_path)
    asset = BuildReferenceAsset(
        reference_id="ref-1",
        file_name="one.jpg",
        source_drive_id="drive-1",
        source_item_id="item-1",
        source_etag="etag-1",
    )
    folder = paths.workspace_reference_cache_dir / "ref-1"
    folder.mkdir(parents=True)
    cached = folder / "one.jpg"
    _photo(cached, (10, 20, 30))
    (folder / "source.json").write_text(json.dumps({"source_etag": "etag-1"}), encoding="utf-8")

    result = resolve_reference_media(asset, paths)

    assert result.path == cached
    assert result.from_cache is True
