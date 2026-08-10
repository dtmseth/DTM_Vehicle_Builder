"""Golden-master digest: canonical, human-diffable normalization of build outputs.

Design spec: docs/audit/GOLDEN_MASTER_SPEC.md (§8.1 Step 1a of the audit roadmap).

Two public entry points:

    pptx_digest(path)  -> dict   normalized slide/shape tree of a .pptx
    json_digest(path)  -> dict   canonicalized parse of a JSON output (plan file)
    canonical_dumps(d) -> str    the one true serialization for storing/diffing

The digest deliberately EXCLUDES everything in the empirical nondeterminism
catalog (see the spec §2): the timestamped output filename, zip entry mtimes,
and docProps package metadata. It also normalizes identity-ish noise that is
deterministic today but carries no visual meaning (rIds, media part names,
numeric shape ids) so mechanical refactors don't produce false diffs.

Everything visual is INCLUDED: slide order, z-order, shape geometry (EMU),
rotation/flips, text runs with formatting, table contents, image content
hashes, fill/line colors.
"""
from __future__ import annotations

import hashlib
import json
import re
from pathlib import Path

from pptx import Presentation
from pptx.oxml.ns import qn


# ── canonical serialization ───────────────────────────────────────────────────

def canonical_dumps(data) -> str:
    """The one serialization used for stored digests and comparisons.

    sort_keys makes the digest insensitive to cosmetic dict-insertion-order
    changes in future code; list order is preserved (it is semantic: slide
    order, z-order, paragraph order, table rows).
    """
    return json.dumps(data, indent=2, sort_keys=True, ensure_ascii=False) + "\n"


def sha256_hex(data: bytes) -> str:
    return hashlib.sha256(data).hexdigest()


# ── JSON outputs (plan file) ──────────────────────────────────────────────────

def json_digest(path: str | Path) -> dict:
    """Parse-and-rewrap a JSON output so byte-level noise can never matter."""
    return json.loads(Path(path).read_text("utf-8"))


# ── PPTX normalization ────────────────────────────────────────────────────────

# Auto-assigned shape names look like "Picture 12" / "Group 5" / "TextBox 3".
# The numeric suffix is an internal insertion id: visually meaningless, and a
# renumbering-only change must not fail the golden master. Z-order (list
# position) already pins ordering.
_AUTO_NAME_RE = re.compile(r"^(.*?)\s+\d+$")


def _normalized_name(name: str) -> str:
    m = _AUTO_NAME_RE.match(name or "")
    return m.group(1) if m else (name or "")


def _color_of(color_format) -> str | None:
    """Best-effort '#RRGGBB' or theme-color name; None when unset/N-A."""
    try:
        ctype = color_format.type
        if ctype is None:
            return None
        try:
            return f"#{color_format.rgb}"
        except Exception:
            return str(color_format.theme_color)
    except Exception:
        return None


def _fill_summary(fill) -> dict | None:
    try:
        ftype = fill.type
    except Exception:
        return None
    if ftype is None:
        return None
    summary: dict = {"type": str(ftype)}
    fore = None
    try:
        fore = _color_of(fill.fore_color)
    except Exception:
        pass
    if fore:
        summary["color"] = fore
    return summary


def _run_summary(run) -> dict:
    font = run.font
    out: dict = {"text": run.text}
    if font.bold is not None:
        out["bold"] = font.bold
    if font.italic is not None:
        out["italic"] = font.italic
    if font.size is not None:
        out["size_pt"] = font.size.pt
    if font.name:
        out["font"] = font.name
    color = _color_of(font.color)
    if color:
        out["color"] = color
    return out


def _text_frame_summary(text_frame) -> list[dict]:
    paragraphs = []
    for para in text_frame.paragraphs:
        p: dict = {"runs": [_run_summary(r) for r in para.runs]}
        if para.alignment is not None:
            p["align"] = str(para.alignment)
        if para.level:
            p["level"] = para.level
        paragraphs.append(p)
    return paragraphs


def _table_summary(table) -> dict:
    return {
        "rows": len(table.rows),
        "cols": len(table.columns),
        "col_widths_emu": [c.width for c in table.columns],
        "row_heights_emu": [r.height for r in table.rows],
        "cells": [
            [
                {
                    "text": cell.text,
                    "paragraphs": _text_frame_summary(cell.text_frame),
                }
                for cell in row.cells
            ]
            for row in table.rows
        ],
    }


def _xfrm_flips(shape) -> tuple[bool, bool]:
    """Read flipH/flipV directly from the XML (not exposed by python-pptx)."""
    try:
        sp_pr = shape._element.spPr
        xfrm = sp_pr.find(qn("a:xfrm")) if sp_pr is not None else None
        if xfrm is None:
            return False, False
        return xfrm.get("flipH") == "1", xfrm.get("flipV") == "1"
    except Exception:
        return False, False


def _shape_summary(shape) -> dict:
    out: dict = {
        "kind": str(shape.shape_type),
        "name": _normalized_name(shape.name),
    }
    for attr in ("left", "top", "width", "height"):
        value = getattr(shape, attr, None)
        if value is not None:
            out[f"{attr}_emu"] = int(value)
    rotation = getattr(shape, "rotation", 0.0)
    if rotation:
        out["rotation"] = rotation
    flip_h, flip_v = _xfrm_flips(shape)
    if flip_h:
        out["flip_h"] = True
    if flip_v:
        out["flip_v"] = True

    # Group: recurse (children in z-order); geometry above is the group box.
    if shape.shape_type is not None and str(shape.shape_type) == "GROUP (6)":
        out["children"] = [_shape_summary(child) for child in shape.shapes]
        return out

    if getattr(shape, "has_text_frame", False) and shape.has_text_frame:
        out["text"] = shape.text_frame.text
        out["paragraphs"] = _text_frame_summary(shape.text_frame)

    if getattr(shape, "has_table", False) and shape.has_table:
        out["table"] = _table_summary(shape.table)

    if shape.shape_type is not None and str(shape.shape_type) == "PICTURE (13)":
        # Identify the image by its bytes, not its rId or media part name.
        try:
            out["image_sha256"] = sha256_hex(shape.image.blob)
            out["image_ext"] = shape.image.ext
        except Exception:
            out["image_sha256"] = None
        try:
            crop = (shape.crop_left, shape.crop_top, shape.crop_right, shape.crop_bottom)
            if any(crop):
                out["crop_ltrb"] = list(crop)
        except Exception:
            pass

    fill = _fill_summary(getattr(shape, "fill", None)) if hasattr(shape, "fill") else None
    if fill:
        out["fill"] = fill
    try:
        line_color = _color_of(shape.line.color)
        if line_color:
            out["line_color"] = line_color
    except Exception:
        pass

    return out


def _slide_summary(slide) -> dict:
    return {
        "layout": slide.slide_layout.name,
        "shapes": [_shape_summary(s) for s in slide.shapes],
    }


def pptx_digest(path: str | Path) -> dict:
    """Normalized structural digest of a rendered .pptx (see module docstring)."""
    prs = Presentation(str(path))
    return {
        "slide_width_emu": prs.slide_width,
        "slide_height_emu": prs.slide_height,
        "slides": [_slide_summary(s) for s in prs.slides],
    }


# ── CLI: python -m tests.golden.digest <file.pptx|file.json> ─────────────────

def main(argv: list[str]) -> int:
    if len(argv) != 2:
        print("usage: python -m tests.golden.digest <output.pptx|output.json>")
        return 2
    target = Path(argv[1])
    if target.suffix.lower() == ".pptx":
        digest = pptx_digest(target)
    else:
        digest = json_digest(target)
    print(canonical_dumps(digest), end="")
    return 0


if __name__ == "__main__":
    import sys

    raise SystemExit(main(sys.argv))
